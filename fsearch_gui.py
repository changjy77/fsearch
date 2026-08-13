#!/usr/bin/env python3
"""
fsearch GUI - PyQt5 기반 파일 검색 도구
"""

import sys
import os
import io
import re
import time
import json
import base64
import subprocess
import logging
import zipfile
import sqlite3
import threading
import multiprocessing
import xml.etree.ElementTree as ET
from datetime import datetime
from pathlib import Path
from concurrent.futures import ThreadPoolExecutor, ProcessPoolExecutor, as_completed
from typing import List, Optional
from collections import defaultdict

from PyQt5.QtWidgets import (
    QApplication, QMainWindow, QWidget, QVBoxLayout, QHBoxLayout, QLayout,
    QLabel, QLineEdit, QPushButton, QTextEdit, QFileDialog, QCheckBox,
    QSpinBox, QProgressBar, QComboBox, QTabWidget, QTableWidget, QTableWidgetItem,
    QSplitter, QMessageBox, QStyledItemDelegate, QFileIconProvider, QMenu
)
from PyQt5.QtCore import (
    Qt, QThread, pyqtSignal, QTimer, QRect, QObject, QEvent, QSize, QPoint,
    QFileInfo, QByteArray, QBuffer, QIODevice
)
from PyQt5.QtGui import QIcon, QColor, QFont, QCursor, QPainter

# 파일 형식별 텍스트 추출 라이브러리
try:
    from docx import Document
except ImportError:
    Document = None

try:
    import pymupdf as fitz
except ImportError:
    fitz = None

try:
    from openpyxl import load_workbook, Workbook
except ImportError:
    load_workbook = None
    Workbook = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    import html2text
except ImportError:
    html2text = None

try:
    import olefile
except ImportError:
    olefile = None


def setup_logging():
    """로깅 설정"""
    log_dir = Path.home() / ".fsearch" / "logs"
    log_dir.mkdir(parents=True, exist_ok=True)

    log_file = log_dir / f"fsearch_{datetime.now().strftime('%Y%m%d')}.log"

    logging.basicConfig(
        level=logging.INFO,
        format='%(asctime)s - %(levelname)s - %(message)s',
        handlers=[
            logging.FileHandler(log_file, encoding='utf-8'),
            logging.StreamHandler()
        ]
    )
    return logging.getLogger(__name__)


class SearchHistory:
    """검색 이력 관리 - 검색어(최대 10개), 경로(최대 5개), 제외 폴더 통계"""
    def __init__(self):
        self.config_dir = Path.home() / ".fsearch"
        self.config_dir.mkdir(parents=True, exist_ok=True)
        self.history_file = self.config_dir / "search_history.json"
        self.data = self._load_data()

    def _load_data(self):
        """저장된 데이터 로드"""
        if self.history_file.exists():
            try:
                with open(self.history_file, 'r', encoding='utf-8') as f:
                    data = json.load(f)
                    # 기존 리스트 형식 호환성 처리
                    if isinstance(data, list):
                        return {
                            'keywords': data,
                            'paths': [],
                            'excluded_stats': {}
                        }
                    elif isinstance(data, dict):
                        # 누락된 키 채우기
                        if 'paths' not in data:
                            data['paths'] = []
                        if 'excluded_stats' not in data:
                            data['excluded_stats'] = {}
                        return data
            except Exception as e:
                logging.getLogger(__name__).warning(f"검색 기록 로드 실패: {e}")
        return {
            'keywords': [],
            'paths': [],
            'excluded_stats': {}
        }

    def _save_data(self):
        """데이터 저장"""
        try:
            with open(self.history_file, 'w', encoding='utf-8') as f:
                json.dump(self.data, f, ensure_ascii=False, indent=2)
        except Exception as e:
            logging.getLogger(__name__).warning(f"검색 기록 저장 실패: {e}")

    def add_keyword(self, keyword):
        """검색어 추가 (중복 제거, 최대 10개)"""
        if not keyword or not keyword.strip():
            return

        keyword = keyword.strip()
        keywords = self.data.get('keywords', [])

        # 중복 제거
        if keyword in keywords:
            keywords.remove(keyword)

        # 맨 앞에 추가
        keywords.insert(0, keyword)

        # 최대 10개 유지
        self.data['keywords'] = keywords[:10]

        self._save_data()

    def add_path(self, path):
        """검색 경로 추가 (중복 제거, 최대 5개)"""
        if not path or not path.strip():
            return

        path = path.strip()
        paths = self.data.get('paths', [])

        # 중복 제거
        if path in paths:
            paths.remove(path)

        # 맨 앞에 추가
        paths.insert(0, path)

        # 최대 5개 유지
        self.data['paths'] = paths[:5]

        self._save_data()

    def add_excluded_folder(self, folder_path):
        """제외 폴더 통계 기록"""
        if not folder_path:
            return

        stats = self.data.get('excluded_stats', {})
        stats[folder_path] = stats.get(folder_path, 0) + 1
        self.data['excluded_stats'] = stats

        self._save_data()

    def clear_keywords(self):
        """모든 검색어 삭제"""
        self.data['keywords'] = []
        self._save_data()

    def clear_paths(self):
        """모든 검색 경로 삭제"""
        self.data['paths'] = []
        self._save_data()

    def clear_excluded_stats(self):
        """자주 제외하는 폴더 통계 삭제"""
        self.data['excluded_stats'] = {}
        self._save_data()

    def get_keywords(self):
        """모든 검색 키워드 반환"""
        return self.data.get('keywords', [])

    def get_paths(self):
        """모든 검색 경로 반환"""
        return self.data.get('paths', [])

    def get_top_excluded_folders(self, limit=5):
        """자주 제외하는 폴더 반환 (상위 N개)"""
        stats = self.data.get('excluded_stats', {})
        if not stats:
            return []
        sorted_items = sorted(stats.items(), key=lambda x: x[1], reverse=True)
        return [folder for folder, _ in sorted_items[:limit]]

    def get_excluded_stats(self):
        """모든 제외 폴더 통계 반환"""
        return self.data.get('excluded_stats', {})


class FlowLayout(QLayout):
    """창 너비에 맞춰 자식 위젯을 자동으로 다음 줄로 줄바꿈하는 레이아웃"""
    def __init__(self, parent=None, margin=0, spacing=6):
        super().__init__(parent)
        if parent is not None:
            self.setContentsMargins(margin, margin, margin, margin)
        self.setSpacing(spacing)
        self._item_list = []

    def addItem(self, item):
        self._item_list.append(item)

    def count(self):
        return len(self._item_list)

    def itemAt(self, index):
        if 0 <= index < len(self._item_list):
            return self._item_list[index]
        return None

    def takeAt(self, index):
        if 0 <= index < len(self._item_list):
            return self._item_list.pop(index)
        return None

    def expandingDirections(self):
        return Qt.Orientations(Qt.Orientation(0))

    def hasHeightForWidth(self):
        return True

    def heightForWidth(self, width):
        return self._do_layout(QRect(0, 0, width, 0), True)

    def setGeometry(self, rect):
        super().setGeometry(rect)
        self._do_layout(rect, False)

    def sizeHint(self):
        return self.minimumSize()

    def minimumSize(self):
        size = QSize()
        for item in self._item_list:
            size = size.expandedTo(item.minimumSize())
        left, top, right, bottom = self.getContentsMargins()
        size += QSize(left + right, top + bottom)
        return size

    def _do_layout(self, rect, test_only):
        left, top, right, bottom = self.getContentsMargins()
        effective_rect = rect.adjusted(left, top, -right, -bottom)
        x = effective_rect.x()
        y = effective_rect.y()
        line_height = 0
        spacing = self.spacing()

        for item in self._item_list:
            next_x = x + item.sizeHint().width() + spacing
            if next_x - spacing > effective_rect.right() and line_height > 0:
                x = effective_rect.x()
                y = y + line_height + spacing
                next_x = x + item.sizeHint().width() + spacing
                line_height = 0

            if not test_only:
                item.setGeometry(QRect(QPoint(x, y), item.sizeHint()))

            x = next_x
            line_height = max(line_height, item.sizeHint().height())

        return y + line_height - rect.y() + bottom


class HoverableTableWidget(QTableWidget):
    """호버 효과가 있는 테이블 위젯"""
    def __init__(self, parent=None):
        super().__init__(parent)
        self.last_hovered_row = -1
        self.hover_callback = None

    def set_hover_callback(self, callback):
        """호버 콜백 설정"""
        self.hover_callback = callback

    def leaveEvent(self, event):
        """마우스가 테이블을 벗어날 때"""
        if self.last_hovered_row >= 0:
            for col in range(self.columnCount()):
                # QTableWidgetItem 처리
                item = self.item(self.last_hovered_row, col)
                if item:
                    item.setBackground(QColor(255, 255, 255))
                # QLabel 위젯 처리
                widget = self.cellWidget(self.last_hovered_row, col)
                if widget and isinstance(widget, QLabel):
                    widget.setStyleSheet("background-color: white;")
            self.last_hovered_row = -1
        super().leaveEvent(event)


class SearchResultDelegate(QStyledItemDelegate):
    """검색 결과에서 검색 단어를 굵게 표시하는 delegate"""

    def __init__(self, keyword, use_regex=False):
        super().__init__()
        self.keyword = keyword
        self.use_regex = use_regex
        if use_regex:
            try:
                self.pattern = re.compile(keyword, re.IGNORECASE)
            except:
                self.pattern = None
        else:
            self.pattern = None

    def paint(self, painter, option, index):
        """검색 단어를 굵게 표시하여 렌더링"""
        text = index.data(Qt.DisplayRole)
        if not text:
            super().paint(painter, option, index)
            return

        # 검색 단어 찾기
        if self.use_regex and self.pattern:
            matches = list(self.pattern.finditer(str(text)))
        else:
            matches = []
            keyword_lower = self.keyword.lower()
            text_lower = str(text).lower()
            start = 0
            while True:
                pos = text_lower.find(keyword_lower, start)
                if pos == -1:
                    break
                matches.append((pos, pos + len(self.keyword)))
                start = pos + 1

        if not matches:
            super().paint(painter, option, index)
            return

        # 배경색 설정
        painter.fillRect(option.rect, option.palette.base())

        # 텍스트와 bold 텍스트 렌더링
        painter.setPen(option.palette.text().color())

        bold_font = QFont(option.font)
        bold_font.setBold(True)

        x = option.rect.left() + 2
        y = option.rect.top() + option.fontMetrics.ascent() + 2

        text_str = str(text)
        last_pos = 0

        for start, end in matches:
            # 일반 텍스트
            if start > last_pos:
                normal_text = text_str[last_pos:start]
                painter.setFont(option.font)
                fm = painter.fontMetrics()
                painter.drawText(x, y, normal_text)
                x += fm.width(normal_text)

            # 굵은 텍스트
            bold_text = text_str[start:end]
            painter.setFont(bold_font)
            fm_bold = painter.fontMetrics()
            painter.drawText(x, y, bold_text)
            x += fm_bold.width(bold_text)

            last_pos = end

        # 남은 텍스트
        if last_pos < len(text_str):
            remaining_text = text_str[last_pos:]
            painter.setFont(option.font)
            fm = painter.fontMetrics()
            painter.drawText(x, y, remaining_text)


def _format_remaining(seconds: float) -> str:
    """초 단위 잔여 시간을 "N분 M초" 형태로 변환 (상태표시바 남은 시간 표기용)"""
    seconds = max(0, int(seconds))
    minutes, sec = divmod(seconds, 60)
    hours, minutes = divmod(minutes, 60)
    if hours:
        return f"{hours}시간 {minutes}분"
    if minutes:
        return f"{minutes}분 {sec}초"
    return f"{sec}초"


def _parse_keyword_query(keyword: str) -> List[List[str]]:
    """검색어를 OR로 구분된 AND-항 그룹들로 파싱한다.
    공백은 AND, 대문자 "OR" 토큰만 명시적 OR 구분자로 인식한다(소문자 "or"는 검색어
    그대로 취급). 예: "네이버 회의록 OR 카카오" -> [["네이버", "회의록"], ["카카오"]]
    반환되는 각 항은 매칭에 바로 쓸 수 있도록 소문자로 변환되어 있다."""
    groups = []
    current = []
    for tok in keyword.split():
        if tok == 'OR':
            if current:
                groups.append(current)
                current = []
        else:
            current.append(tok.lower())
    if current:
        groups.append(current)
    if not groups:
        groups = [[keyword.strip().lower()]]
    return groups


def _query_matches(query: List[List[str]], text_lower: str) -> bool:
    """OR-그룹 중 하나라도 모든 AND-항이 text_lower에 포함되면 True"""
    return any(all(term in text_lower for term in group) for group in query)


def _query_any_term_in(query: List[List[str]], text_lower: str) -> bool:
    """쿼리를 구성하는 개별 항 중 하나라도 text_lower에 있으면 True
    (라인 단위 미리보기/하이라이트처럼 AND 전체 충족 여부가 아니라 부분 일치 표시가 필요할 때 사용)"""
    return any(term in text_lower for group in query for term in group)


def _can_use_index(keyword: str) -> bool:
    """이 키워드를 trigram 인덱스로 안전하고 빠르게 찾을 수 있는지 판단.

    - 3글자 미만은 trigram 인덱스를 쓰지 못해 전체 스캔이 되고, 오히려 기존 방식보다 느리다.
    - SQLite의 LIKE 대소문자 무시는 ASCII 전용이라 'café'로 'CAFÉ'를 찾지 못한다.
      파이썬의 lower()와 결과가 달라지므로(누락 발생) 그런 키워드는 인덱스를 쓰지 않는다.
      한글·한자처럼 대소문자가 없는 문자는 영향이 없어 빠른 경로를 그대로 탄다.
    """
    if len(keyword) < 3:
        return False
    return not any(ord(c) > 127 and c.lower() != c.upper() for c in keyword)


# 파일 형식별 추출 로직 버전.
#
# 캐시는 (mtime, size)로만 유효성을 판단하므로, 추출 코드를 고쳐도 파일 자체가
# 그대로면 옛 추출 결과를 계속 쓴다. 실제로 .hwp 본문 추출과 CP949 텍스트,
# .docx 표/텍스트박스를 고친 뒤에도 이미 캐시돼 있던 파일 147개가 검색에서
# 계속 누락되고 있었고, 캐시를 통째로 비우고 나서야 드러났다.
#
# **해당 형식의 추출 코드를 고치면 여기 번호를 올린다.** 번호가 바뀐 형식만
# 캐시에서 지워져 다음 검색 때 다시 추출된다(나머지 형식은 그대로 재사용).
EXTRACT_VERSIONS = {
    '.txt': 1, '.md': 1, '.csv': 1, '.json': 1, '.xml': 1,
    '.html': 1, '.htm': 1,
    '.docx': 1, '.pptx': 1, '.pdf': 1,
    '.xlsx': 1, '.xls': 1,
    '.hwp': 1, '.hwpx': 1,
}


class TextExtractionCache:
    """파일에서 추출한 텍스트를 디스크(SQLite)에 캐싱 - 재검색 시 파싱 재사용
    본문 검색은 trigram FTS5 인덱스로 처리해 전체 본문을 메모리로 올리지 않는다."""

    def __init__(self):
        cache_dir = Path.home() / ".fsearch"
        cache_dir.mkdir(parents=True, exist_ok=True)
        self.db_path = cache_dir / "text_cache.db"
        conn = sqlite3.connect(str(self.db_path))
        conn.execute(
            "CREATE TABLE IF NOT EXISTS text_cache "
            "(path TEXT PRIMARY KEY, mtime REAL, size INTEGER, text TEXT)"
        )
        conn.execute("CREATE TABLE IF NOT EXISTS cache_meta (key TEXT PRIMARY KEY, value TEXT)")
        # trigram 토크나이저라야 '네이버'가 '네이버클라우드'에 매칭되는 부분 문자열 검색이 된다
        # (기본 unicode61은 단어 단위라 기존 검색 동작과 결과가 달라짐)
        # trigram은 SQLite 3.34+ 에서만 지원되므로, 낮은 버전에서는 인덱스 없이 동작하도록 한다.
        # 인덱스가 없어도 캐시 본문 전체를 훑는 기존 경로로 검색되며 결과는 동일하다(속도만 느려짐).
        try:
            conn.execute(
                "CREATE VIRTUAL TABLE IF NOT EXISTS text_fts USING fts5("
                "text, content='text_cache', content_rowid='rowid', tokenize='trigram')"
            )
            self._create_triggers(conn)
            self.fts_available = True
        except sqlite3.DatabaseError:
            self.fts_available = False

        # 새 캐시(빈 DB)는 트리거만으로 인덱스가 항상 최신이므로 재구축이 필요 없다
        if self.fts_available and not conn.execute("SELECT EXISTS(SELECT 1 FROM text_cache)").fetchone()[0]:
            conn.execute("INSERT OR REPLACE INTO cache_meta (key, value) VALUES ('fts_built', '1')")

        # 추출 로직이 바뀐 형식의 캐시 정리 (트리거가 있어야 FTS 인덱스도 함께 정리된다)
        self.cleared_by_version = self._apply_extract_versions(conn)
        conn.commit()
        conn.close()

    @staticmethod
    def _apply_extract_versions(conn):
        """EXTRACT_VERSIONS와 캐시에 기록된 버전을 비교해, 바뀐 형식의 캐시만 지운다.
        지워진 항목은 다음 검색에서 새 추출 로직으로 다시 채워진다."""
        stored = dict(conn.execute(
            "SELECT key, value FROM cache_meta WHERE key LIKE 'extver:%'"
        ).fetchall())

        has_rows = conn.execute("SELECT EXISTS(SELECT 1 FROM text_cache)").fetchone()[0]
        if not stored and has_rows:
            # 버전 기록이 없는 기존 캐시는 현재 로직으로 만들어진 것으로 보고 표시만 남긴다.
            # 여기서 전부 지우면 수 GB 캐시를 통째로 다시 만들게 되어 손해가 더 크다.
            conn.executemany(
                "INSERT OR REPLACE INTO cache_meta (key, value) VALUES (?, ?)",
                [(f'extver:{ext}', str(ver)) for ext, ver in EXTRACT_VERSIONS.items()]
            )
            return 0

        cleared = 0
        for ext, ver in EXTRACT_VERSIONS.items():
            key = f'extver:{ext}'
            if stored.get(key) == str(ver):
                continue
            # zip 내부 항목은 'zip경로!내부파일명' 형태라 확장자로 끝나는 것은 동일하다
            cur = conn.execute("DELETE FROM text_cache WHERE LOWER(path) LIKE ?", (f'%{ext}',))
            cleared += cur.rowcount
            conn.execute(
                "INSERT OR REPLACE INTO cache_meta (key, value) VALUES (?, ?)", (key, str(ver))
            )

        if cleared:
            # zip 스캔 표시를 지워 내부 항목을 다시 열거하게 한다.
            # (표시가 남아 있으면 방금 지운 내부 항목을 다시 추출하지 않는다)
            conn.execute(
                "DELETE FROM text_cache WHERE LOWER(path) LIKE '%.zip' AND path NOT LIKE '%!%'"
            )
        return cleared

    TRIGGER_NAMES = ('text_cache_ai', 'text_cache_ad', 'text_cache_au')

    @classmethod
    def _create_triggers(cls, conn):
        """external content 방식은 자동 동기화되지 않으므로 트리거로 본문 변경을 인덱스에 반영"""
        conn.execute(
            "CREATE TRIGGER IF NOT EXISTS text_cache_ai AFTER INSERT ON text_cache BEGIN "
            "INSERT INTO text_fts(rowid, text) VALUES (new.rowid, new.text); END"
        )
        conn.execute(
            "CREATE TRIGGER IF NOT EXISTS text_cache_ad AFTER DELETE ON text_cache BEGIN "
            "INSERT INTO text_fts(text_fts, rowid, text) VALUES('delete', old.rowid, old.text); END"
        )
        conn.execute(
            "CREATE TRIGGER IF NOT EXISTS text_cache_au AFTER UPDATE ON text_cache BEGIN "
            "INSERT INTO text_fts(text_fts, rowid, text) VALUES('delete', old.rowid, old.text); "
            "INSERT INTO text_fts(rowid, text) VALUES (new.rowid, new.text); END"
        )

    @staticmethod
    def _install_stop_handler(conn, stop_check):
        """stop_check()가 True를 반환하면 실행 중인 SQL을 중단시킨다(SQLITE_INTERRUPT).
        1000 VM 명령마다 호출되며 오버헤드는 무시할 수준이다(실측: 유무 차이 0.01초 이내).
        rebuild/VACUUM은 원자적이라 중단돼도 트랜잭션이 롤백되어 이전 상태로 안전하게 남는다."""
        if stop_check is None:
            return
        conn.set_progress_handler(lambda: 1 if stop_check() else 0, 1000)

    def index_needs_build(self) -> bool:
        """기존 캐시가 아직 색인되지 않았는지 확인 (업그레이드 후 최초 1회만 참)"""
        if not self.fts_available:
            return False
        conn = sqlite3.connect(str(self.db_path))
        built = conn.execute("SELECT value FROM cache_meta WHERE key = 'fts_built'").fetchone()
        has_rows = conn.execute("SELECT EXISTS(SELECT 1 FROM text_cache)").fetchone()[0]
        conn.close()
        return not built and bool(has_rows)

    def build_index(self, stop_check=None):
        """기존 캐시 전체를 trigram 인덱스로 색인 (1회성, 데이터량에 비례해 수 분 소요).
        stop_check가 주어지고 도중 True가 되면 rebuild를 중단한다. 중단되면 트랜잭션이
        롤백되어 트리거도 색인도 없는 이전 상태로 남고, 다음 검색에서 다시 시도된다."""
        conn = sqlite3.connect(str(self.db_path))
        # 정리 단계의 UPDATE가 트리거를 통해 아직 비어 있는 인덱스에 delete를 시도하면
        # 인덱스가 손상되므로("database disk image is malformed"), 트리거를 잠시 걷어내고
        # 본문 정리를 끝낸 뒤 rebuild로 인덱스를 통째로 다시 만든다.
        for name in self.TRIGGER_NAMES:
            conn.execute(f"DROP TRIGGER IF EXISTS {name}")
        self._sanitize_existing(conn)
        self._install_stop_handler(conn, stop_check)
        try:
            conn.execute("INSERT INTO text_fts(text_fts) VALUES('rebuild')")
        except sqlite3.OperationalError:
            conn.close()
            return
        conn.execute("INSERT OR REPLACE INTO cache_meta (key, value) VALUES ('fts_built', '1')")
        self._create_triggers(conn)
        conn.commit()
        conn.close()

    # 빈 공간이 이 비율을 넘으면 정리한다. 정리 자체가 수십 초 걸리므로 자주 하지 않는다.
    MAINTENANCE_FREE_RATIO = 0.10
    MAINTENANCE_MIN_SIZE = 200 * 1024 * 1024

    def needs_maintenance(self) -> bool:
        """회수할 빈 공간이 충분히 쌓였는지 확인.
        증분 쓰기가 쌓이면 삭제된 행이 남긴 페이지가 파일에 그대로 남는다
        (실측: 4.26GB 중 658MB)."""
        try:
            size = self.db_path.stat().st_size
        except OSError:
            return False
        if size < self.MAINTENANCE_MIN_SIZE:
            return False
        try:
            conn = sqlite3.connect(str(self.db_path))
            page_size = conn.execute("PRAGMA page_size").fetchone()[0]
            free_pages = conn.execute("PRAGMA freelist_count").fetchone()[0]
            conn.close()
        except sqlite3.DatabaseError:
            return False
        return free_pages * page_size > size * self.MAINTENANCE_FREE_RATIO

    def run_maintenance(self, stop_check=None):
        """FTS 인덱스 조각을 병합하고 빈 공간을 파일에서 회수한다.
        실측(4.26GB / 39,642건): optimize 29초 + VACUUM 62초 -> 3.05GB,
        키워드 조회 17ms -> 4ms.
        stop_check가 주어지고 도중 True가 되면 중단한다. VACUUM은 임시 파일에 새로
        만든 뒤 원본과 교체하는 방식이라, 중단돼도 원본 파일은 그대로 안전하게 남는다
        (실측: integrity_check 'ok' 유지)."""
        # VACUUM은 트랜잭션 안에서 실행할 수 없어 자동 커밋 모드로 연결한다
        conn = sqlite3.connect(str(self.db_path), isolation_level=None)
        self._install_stop_handler(conn, stop_check)
        try:
            if self.fts_available:
                try:
                    conn.execute("INSERT INTO text_fts(text_fts) VALUES('optimize')")
                except sqlite3.DatabaseError:
                    pass
            if stop_check is None or not stop_check():
                conn.execute("VACUUM")
        except sqlite3.OperationalError:
            pass
        finally:
            conn.close()

    def row_count(self) -> int:
        """캐시에 들어 있는 항목 수 (색인 전략 판단용)"""
        conn = sqlite3.connect(str(self.db_path))
        n = conn.execute("SELECT COUNT(*) FROM text_cache").fetchone()[0]
        conn.close()
        return n

    def should_defer_index(self, new_entries: int) -> bool:
        """색인을 나중에 한 번에 만드는 편이 유리한지 판단.

        행마다 트리거로 갱신하면 본문 1M자당 0.194초, 마지막에 rebuild로 한 번에
        만들면 0.110초다(실측). 다만 rebuild는 새 항목만이 아니라 캐시 '전체'를
        다시 색인하므로, 새로 넣는 양이 기존 캐시에 비해 작으면 오히려 손해다.
        여유를 두어 새 항목이 기존 항목 수 이상일 때만 미룬다.
        """
        if not self.fts_available or new_entries < 1000:
            return False
        return new_entries >= self.row_count()

    def begin_deferred_index(self):
        """대량 추출 동안 색인 갱신을 멈춘다.
        도중에 중단되면 fts_built 표시가 없는 상태로 남아, 다음 검색에서
        build_index()가 인덱스와 트리거를 함께 복구한다."""
        conn = sqlite3.connect(str(self.db_path))
        conn.execute("DELETE FROM cache_meta WHERE key = 'fts_built'")
        for name in self.TRIGGER_NAMES:
            conn.execute(f"DROP TRIGGER IF EXISTS {name}")
        conn.commit()
        conn.close()

    def end_deferred_index(self, stop_check=None):
        """미뤄둔 색인을 한 번에 만들고 트리거를 복구한다.
        추출 시점에 이미 NUL을 제거하므로 build_index()와 달리 본문 정리는 하지 않는다.
        stop_check가 주어지고 도중 True가 되면 중단한다. 중단되면 트리거가 없는 채로
        남지만(begin_deferred_index에서 이미 지워둔 상태), 다음 검색에서
        index_needs_build()가 이를 감지해 build_index()로 복구한다."""
        conn = sqlite3.connect(str(self.db_path))
        self._install_stop_handler(conn, stop_check)
        try:
            conn.execute("INSERT INTO text_fts(text_fts) VALUES('rebuild')")
        except sqlite3.OperationalError:
            conn.close()
            return
        conn.execute("INSERT OR REPLACE INTO cache_meta (key, value) VALUES ('fts_built', '1')")
        self._create_triggers(conn)
        conn.commit()
        conn.close()

    @staticmethod
    def _sanitize_existing(conn):
        """기존 캐시에 남아 있는 NUL 문자를 공백으로 치환.
        SQL의 replace()는 NUL에서 잘려 쓸 수 없으므로 파이썬에서 바이트 단위로 처리한다."""
        fixes = []
        for rowid, blob in conn.execute(
            "SELECT rowid, CAST(text AS BLOB) FROM text_cache WHERE text IS NOT NULL"
        ):
            if blob and b'\x00' in blob:
                fixes.append((blob.replace(b'\x00', b' ').decode('utf-8', errors='ignore'), rowid))
        if fixes:
            conn.executemany("UPDATE text_cache SET text = ? WHERE rowid = ?", fixes)
        return len(fixes)

    def load_for_prefix(self, path_prefix):
        """지정 경로 하위의 캐시만 로드하여 {경로: (mtime, size, text)} 반환"""
        conn = sqlite3.connect(str(self.db_path))
        rows = conn.execute(
            "SELECT path, mtime, size, text FROM text_cache WHERE path LIKE ?",
            (self._prefix_pattern(path_prefix),)
        ).fetchall()
        conn.close()
        return {path: (mtime, size, text) for path, mtime, size, text in rows}

    def load_meta_for_prefix(self, path_prefix):
        """본문 없이 {경로: (mtime, size)}만 로드 - 캐시 유효성 판단과 정리에는 본문이 필요 없다"""
        conn = sqlite3.connect(str(self.db_path))
        rows = conn.execute(
            "SELECT path, mtime, size FROM text_cache WHERE path LIKE ?",
            (self._prefix_pattern(path_prefix),)
        ).fetchall()
        conn.close()
        return {path: (mtime, size) for path, mtime, size in rows}

    def load_matching_texts(self, path_prefix, keyword):
        """캐시된 본문 중 키워드를 포함하는 항목만 {경로: 텍스트}로 반환 (trigram 인덱스 사용)

        ESCAPE 절을 붙이면 trigram 인덱스 최적화가 비활성화된다(실측 13ms -> 1755ms).
        그래서 평소에는 이스케이프 없이 조회한다. 키워드의 %와 _가 와일드카드로 동작하지만
        결과는 항상 정답의 상위집합이므로 가져온 본문으로 파이썬에서 정확히 걸러낸다.
        다만 키워드에 실제로 와일드카드가 들어 있으면 상위집합이 지나치게 커져 오히려 느려지므로
        (실측 '50%' 4891ms) 그때는 이스케이프해서 조회한다.
        """
        conn = sqlite3.connect(str(self.db_path))
        if '%' in keyword or '_' in keyword:
            pattern = keyword.replace('\\', '\\\\').replace('%', '\\%').replace('_', '\\_')
            sql = ("SELECT path, text FROM text_cache WHERE rowid IN "
                   "(SELECT rowid FROM text_fts WHERE text LIKE ? ESCAPE '\\') AND path LIKE ?")
        else:
            pattern = keyword
            sql = ("SELECT path, text FROM text_cache WHERE rowid IN "
                   "(SELECT rowid FROM text_fts WHERE text LIKE ?) AND path LIKE ?")
        rows = conn.execute(sql, (f"%{pattern}%", self._prefix_pattern(path_prefix))).fetchall()
        conn.close()
        keyword_lower = keyword.lower()
        return {path: text for path, text in rows if text and keyword_lower in text.lower()}

    @staticmethod
    def _prefix_pattern(path_prefix):
        """캐시 키는 str(Path(...))로 저장되어 구분자가 '\'이므로 조회 prefix도 동일하게 정규화
        (GUI가 'D:/클로드'처럼 '/'로 넘기면 매칭이 전부 실패해 캐시가 무효화됨)"""
        return str(Path(path_prefix)).rstrip("\\/") + "%"

    def save_entries(self, entries):
        """{경로: (mtime, size, text)} 항목들을 일괄 저장"""
        if not entries:
            return
        conn = sqlite3.connect(str(self.db_path))
        # INSERT OR REPLACE의 암묵적 삭제는 재귀 트리거가 꺼져 있으면 DELETE 트리거를 발생시키지 않아
        # 인덱스에 옛 본문이 남는다. 삭제와 삽입을 명시적으로 나눠 트리거가 확실히 동작하게 한다.
        conn.executemany("DELETE FROM text_cache WHERE path = ?", [(path,) for path in entries])
        conn.executemany(
            "INSERT INTO text_cache (path, mtime, size, text) VALUES (?, ?, ?, ?)",
            [(path, mtime, size, text) for path, (mtime, size, text) in entries.items()]
        )
        conn.commit()
        conn.close()

    def delete_entries(self, keys):
        """지정된 캐시 키(경로 또는 zip경로!내부파일명)들을 일괄 삭제"""
        if not keys:
            return
        conn = sqlite3.connect(str(self.db_path))
        conn.executemany("DELETE FROM text_cache WHERE path = ?", [(k,) for k in keys])
        conn.commit()
        conn.close()


def _as_stream(source):
    """추출 대상을 파서에 넘길 형태로 변환.
    대부분의 파서(python-docx/pptx, openpyxl, zipfile, olefile)는 경로와 파일 객체를
    모두 받으므로, 메모리에 있는 내용은 BytesIO로 감싸 임시파일 없이 넘긴다.
    호출할 때마다 새 객체를 만들어야 한다(한 번 읽은 스트림은 재사용할 수 없다)."""
    return io.BytesIO(source) if isinstance(source, bytes) else source


def _read_source_bytes(source) -> bytes:
    """추출 대상의 원본 바이트"""
    return source if isinstance(source, bytes) else source.read_bytes()


def _read_text_smart(source) -> str:
    """UTF-8로 우선 시도하고, 실패하면 CP949(EUC-KR)로 재시도 (오래된 한글 문서 대응)"""
    raw = _read_source_bytes(source)
    try:
        return raw.decode('utf-8-sig')
    except UnicodeDecodeError:
        return raw.decode('cp949', errors='ignore')


def _extract_docx_raw_xml(source) -> str:
    """word/document.xml에서 텍스트 태그를 직접 추출 (python-docx 파싱 실패 시 우회용)"""
    try:
        with zipfile.ZipFile(_as_stream(source)) as z:
            xml = z.read('word/document.xml').decode('utf-8', errors='ignore')
        paragraphs = re.findall(r'<w:p(?:\s[^>]*)?>.*?</w:p>', xml, re.DOTALL)
        lines = [''.join(re.findall(r'<w:t[^>]*>([^<]*)</w:t>', p)) for p in paragraphs]
        return '\n'.join(lines)
    except:
        return ""


def _extract_docx_textboxes(source) -> str:
    """텍스트박스(w:txbxContent) 안의 텍스트 추출 (python-docx가 지원하지 않아 raw XML로 보완)"""
    try:
        with zipfile.ZipFile(_as_stream(source)) as z:
            xml = z.read('word/document.xml').decode('utf-8', errors='ignore')
        boxes = re.findall(r'<w:txbxContent[^>]*>.*?</w:txbxContent>', xml, re.DOTALL)
        lines = [''.join(re.findall(r'<w:t[^>]*>([^<]*)</w:t>', box)) for box in boxes]
        return '\n'.join(lines)
    except:
        return ""


def _strip_nul(text: str) -> str:
    """NUL 문자를 공백으로 치환.
    SQLite의 문자열 함수와 FTS5 인덱스는 NUL에서 문자열이 끝난 것으로 취급해
    그 뒤 내용이 통째로 검색에서 누락된다. PDF 추출기가 공백 자리에 NUL을 넣는
    경우가 많아(예: '접수:\\x002020.\\x009.') 삭제가 아니라 공백으로 바꾼다."""
    return text.replace('\x00', ' ') if text else text


def _extract_text_impl(file_path: Path) -> Optional[str]:
    """파일 형식별 텍스트 추출 (실제 파싱 로직)
    모듈 레벨 함수 - SearchWorker(QThread)의 인스턴스 메서드는 pickle이 불가능해
    ProcessPoolExecutor 워커로 전달할 수 없으므로 클래스 밖에 둔다."""
    return _strip_nul(_extract_text_raw(file_path, file_path.suffix.lower()))


def _extract_text_from_bytes(data: bytes, ext: str) -> Optional[str]:
    """이미 메모리에 있는 내용에서 바로 텍스트 추출 (zip 내부 파일용).
    임시파일로 썼다가 파서가 다시 읽는 왕복을 없앤다(실측 엔트리당 2.16ms)."""
    return _strip_nul(_extract_text_raw(data, ext))


def _extract_text_raw(source, ext: str):
    """파일 형식별 텍스트 추출 (형식별 분기 본체).
    source는 파일 경로(Path) 또는 파일 내용(bytes) 둘 다 받는다.
    반환값은 성공 시 문자열(빈 문서면 ""), 파싱 자체가 실패하면 None이다.
    이 구분이 있어야 호출부(SearchWorker.run())가 "정말 실패한 파일"을
    집계해 로그로 남길 수 있다 - 이전에는 둘 다 ""라 조용히 묻혔다."""
    try:
        if ext == '.txt' or ext == '.md' or ext == '.csv' or ext == '.json' or ext == '.xml':
            # 텍스트/마크다운/CSV/JSON/XML 파일
            return _read_text_smart(source)

        elif ext == '.html' or ext == '.htm':
            # HTML 파일
            return _read_text_smart(source)

        elif ext == '.docx' and Document:
            # Word 문서 - 본문 외에 표/머리글/바닥글/텍스트박스도 포함
            try:
                doc = Document(_as_stream(source))
                parts = [para.text for para in doc.paragraphs]

                for table in doc.tables:
                    for row in table.rows:
                        for cell in row.cells:
                            parts.extend(p.text for p in cell.paragraphs)

                for section in doc.sections:
                    parts.extend(p.text for p in section.header.paragraphs)
                    parts.extend(p.text for p in section.footer.paragraphs)

                textbox_text = _extract_docx_textboxes(source)
                if textbox_text:
                    parts.append(textbox_text)

                return '\n'.join(parts)
            except:
                # 임베디드 첨부 등으로 python-docx 파싱 실패 시 document.xml 직접 추출로 우회
                return _extract_docx_raw_xml(source)

        elif ext == '.pptx' and Presentation:
            # 파워포인트 문서
            try:
                prs = Presentation(_as_stream(source))
                text = ""
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            text += shape.text_frame.text + "\n"
                return text
            except Exception:
                return None

        elif ext == '.pdf' and fitz:
            # PDF 파일 (PyMuPDF만 경로/스트림 인자가 달라 따로 분기한다)
            try:
                text = ""
                if isinstance(source, bytes):
                    doc = fitz.open(stream=source, filetype='pdf')
                else:
                    doc = fitz.open(str(source))
                with doc:
                    for page in doc:
                        text += page.get_text() + "\n"
                return text
            except Exception:
                return None

        elif ext in ['.xlsx', '.xls'] and load_workbook:
            # Excel 파일
            try:
                if ext == '.xlsx':
                    wb = load_workbook(_as_stream(source), data_only=True)
                    text = ""
                    for ws in wb.sheetnames:
                        sheet = wb[ws]
                        for row in sheet.iter_rows():
                            for cell in row:
                                if cell.value:
                                    text += str(cell.value) + " "
                            text += "\n"
                    return text
                else:
                    # .xls 파일은 openpyxl로 지원하지 않음
                    import xlrd
                    if isinstance(source, bytes):
                        wb = xlrd.open_workbook(file_contents=source)
                    else:
                        wb = xlrd.open_workbook(str(source))
                    text = ""
                    for sheet in wb.sheets():
                        for row in sheet.get_rows():
                            for cell in row:
                                text += str(cell.value) + " "
                            text += "\n"
                    return text
            except Exception:
                return None

        elif ext == '.hwpx':
            # 한글 파일 (.hwpx, zip 기반) - Contents/section*.xml에서 <hp:t> 텍스트 추출
            try:
                text = ""
                with zipfile.ZipFile(_as_stream(source), 'r') as hwp:
                    for name in hwp.namelist():
                        if 'section' in name.lower() and name.lower().endswith('.xml'):
                            try:
                                root = ET.fromstring(hwp.read(name))
                                for elem in root.iter():
                                    if (elem.tag == 't' or elem.tag.endswith('}t')) and elem.text:
                                        text += elem.text + " "
                            except ET.ParseError:
                                pass
                return text if text else ""
            except Exception:
                return None

        elif ext == '.hwp':
            # 구버전 한글 파일 (OLE2 복합문서) - PrvText(미리보기 텍스트) 스트림에서 추출
            if olefile is None:
                return ""
            try:
                with olefile.OleFileIO(_as_stream(source)) as ole:
                    if not ole.exists('PrvText'):
                        return ""
                    data = ole.openstream('PrvText').read()
                return data.decode('utf-16le', errors='ignore')
            except Exception:
                return None

        else:
            return ""

    except Exception:
        return None


def _extract_text_worker(file_path_str: str):
    """ProcessPoolExecutor 워커 진입점: 파일 경로 -> (경로, mtime, size, 추출된 텍스트).
    CPU 바운드 파싱(PDF/docx/pptx/xlsx/hwp)은 GIL 때문에 스레드로는 병렬성이 거의 나오지 않아
    (실측: 16스레드로도 1.03배) 별도 프로세스로 실행한다."""
    file_path = Path(file_path_str)
    try:
        stat = file_path.stat()
        mtime, size = stat.st_mtime, stat.st_size
    except OSError:
        return (file_path_str, None, None, "")
    text = _extract_text_impl(file_path)
    return (file_path_str, mtime, size, text)


def _extract_zip_entries_worker(task):
    """ProcessPoolExecutor 워커 진입점: zip 하나에서 지정된 내부 파일들의 텍스트를 추출.
    zip 내부 파일도 일반 문서와 같은 CPU 바운드 파싱이라 스레드에서는 GIL에 묶여 직렬화된다
    (실측: 엔트리당 23.8ms, 3만여 건 기준 약 13분). zip을 엔트리마다 여는 낭비를 피하려고
    한 번 연 zip에서 여러 엔트리를 처리하고, 결과를 묶어서 돌려준다.

    task: (zip 경로 문자열, [(내부 파일명, 확장자), ...])
    반환: [(캐시키, CRC, 크기, 텍스트), ...]  - 캐시키는 'zip경로!내부파일명'
    """
    zip_path_str, entries = task
    results = []
    try:
        with zipfile.ZipFile(zip_path_str) as zf:
            for inner_name, ext in entries:
                try:
                    info = zf.getinfo(inner_name)
                    data = zf.read(inner_name)
                    text = _extract_text_from_bytes(data, ext)
                    results.append((f"{zip_path_str}!{inner_name}", info.CRC, info.file_size, text))
                except Exception:
                    continue
    except Exception:
        pass
    return results


class SearchWorker(QThread):
    """검색을 별도 스레드에서 실행"""
    progress = pyqtSignal(int)
    result_found = pyqtSignal(dict)  # 결과를 실시간으로 전송
    finished = pyqtSignal(list)
    status = pyqtSignal(str)
    error = pyqtSignal(str)
    excluded_files_updated = pyqtSignal(list)  # 제외된 파일 목록 업데이트
    file_processing = pyqtSignal(str)  # 현재 처리 중인 파일 이름
    no_match_files_updated = pyqtSignal(list)  # 검색어 미포함 파일 목록
    skipped_files_count = pyqtSignal(int)  # 스킵된 파일 갯수
    skipped_files_updated = pyqtSignal(list)  # 스킵된 대용량파일 목록
    extraction_urgent = pyqtSignal(bool)  # 캐시DB가 비어 있어 대량 추출+색인 생성 중임을 표시

    # 내용 추출 가능한 확장자 (zip 내부 검색 시에도 동일하게 사용)
    EXTRACTABLE_EXTENSIONS = {
        '.doc', '.docx', '.ppt', '.pptx', '.hwp', '.hwpx', '.pdf',
        '.xls', '.xlsx', '.txt', '.html', '.htm', '.md', '.csv', '.json', '.xml'
    }

    # 결과가 이보다 많으면 검색을 중단한다. ".doc"처럼 흔한 키워드는 trigram 인덱스가
    # 걸러주지 못해(실측: 캐시 39,644건 중 72%가 매치) 결과가 3만 건 넘게 쏟아지고,
    # 그 결과를 테이블에 실시간 추가+정렬+굵게강조+컬럼폭 재계산하는 GUI 후처리가
    # 결과 수에 비례해 느려져 실측 88초간 응답 없음 상태가 된다.
    MAX_RESULTS = 2000

    def __init__(self, keyword, path, ignore_dirs, name_only, content_only, use_regex, max_workers, skip_large_files=False):
        super().__init__()
        self.keyword = keyword
        # OR로 구분된 AND-항 그룹들. 정규식 모드에서는 self.keyword를 그대로 패턴으로 쓰므로 안 쓴다.
        self.keyword_query = _parse_keyword_query(keyword)
        self.path = path
        self.ignore_dirs = set(ignore_dirs)
        self.name_only = name_only
        self.content_only = content_only
        self.use_regex = use_regex
        self.max_workers = max_workers
        self.skip_large_files = skip_large_files
        self.results = []
        self.excluded_files = []  # 제외된 파일 목록
        self.text_cache = TextExtractionCache()  # 디스크 캐시 관리자
        self.file_cache = {}  # 메모리 캐싱: 파일 경로 → (mtime, size, 텍스트)
        self.new_cache_entries = {}  # 이번 검색에서 새로 추출된 항목 (디스크 저장용)
        self.cache_lock = threading.Lock()  # 캐시 동시 접근 보호
        self.skipped_large_files = 0  # 스킵된 대용량 파일 갯수
        self.skipped_files_list = []  # 스킵된 대용량파일 목록
        self.extract_failed_files = []  # 텍스트 추출(파싱) 자체가 실패한 파일 목록
        self.stop_flag = False  # 검색 중단 플래그
        self._last_file_signal = 0.0  # file_processing 시그널 발행 간격 조절용
        self._found_lock = threading.Lock()  # _found_count 동시 접근 보호
        self._found_count = 0  # 지금까지 찾은 결과 수(스레드 실시간 집계, MAX_RESULTS 판단용)
        self._max_results_hit = False  # stop_flag가 사용자 중지가 아니라 MAX_RESULTS 초과로 세워졌는지

    def run(self):
        """검색 실행"""
        self.timing = {}
        t_run_start = time.time()
        try:
            # 파일 수집
            self.status.emit("📂 파일을 수집 중입니다...")
            files = self._collect_files()
            self.timing['collect'] = time.time() - t_run_start

            if not files:
                self.status.emit("❌ 검색할 파일이 없습니다.")
                self.finished.emit([])
                return

            self.status.emit(f"✅ {len(files)}개 파일 발견. 검색 중...")

            # 정규식 컴파일
            if self.use_regex:
                try:
                    regex = re.compile(self.keyword, re.IGNORECASE)
                except re.error as e:
                    self.error.emit(f"정규식 오류: {e}")
                    return
            else:
                regex = None

            # 디스크 캐시 로드 (같은 폴더 재검색 시 텍스트 추출 재사용)
            # 파일명만 검색할 때는 텍스트 추출 자체를 하지 않으므로 캐시 로드를 생략
            t_phase_start = time.time()
            if not self.name_only:
                if self.text_cache.cleared_by_version:
                    # 이 검색이 평소보다 느린 이유를 알 수 있게 알린다
                    self.status.emit(
                        f"🔄 추출 방식이 바뀐 {self.text_cache.cleared_by_version}개 항목을 다시 읽습니다..."
                    )
                if self.text_cache.index_needs_build():
                    # 기존 캐시를 처음 색인하는 경우이거나, 대량 추출이 도중에 끊겨
                    # 색인이 최신이 아닌 경우다. 어느 쪽이든 여기서 통째로 다시 만든다.
                    self.status.emit("🔧 검색 색인을 만드는 중입니다. 수 분 걸릴 수 있습니다...")
                    self.text_cache.build_index(stop_check=lambda: self.stop_flag)
                    if self.stop_flag:
                        self.status.emit("⏹️ 검색 중단됨")
                        self.finished.emit([])
                        return

                # 캐시 정리는 검색이 끝난 뒤가 아니라 시작 전에 한다.
                # 끝난 뒤에 하면 사용자가 곧바로 다음 검색을 시작했을 때 DB가 잠겨 충돌한다.
                if self.text_cache.needs_maintenance():
                    self.status.emit("🧹 검색 캐시를 정리하는 중입니다. 1~2분 걸릴 수 있습니다...")
                    self.text_cache.run_maintenance(stop_check=lambda: self.stop_flag)
                    if self.stop_flag:
                        self.status.emit("⏹️ 검색 중단됨")
                        self.finished.emit([])
                        return

                is_single_term = len(self.keyword_query) == 1 and len(self.keyword_query[0]) == 1
                if (self.use_regex or not self.text_cache.fts_available
                        or not is_single_term or not _can_use_index(self.keyword)):
                    # 정규식과 AND/OR 다중 키워드는 인덱스로 평가할 수 없고, 짧거나 비ASCII
                    # 대소문자가 섞인 키워드는 인덱스 결과가 파이썬 매칭과 달라질 수 있어
                    # 기존처럼 본문 전체를 로드한다
                    # (SQLite 버전이 낮아 인덱스를 만들지 못한 경우에도 이 경로를 탄다)
                    self.file_cache = self.text_cache.load_for_prefix(self.path)
                else:
                    # 키워드를 포함하지 않는 캐시 항목은 본문을 로드하지 않는다.
                    # 내용 매칭 결과가 0인 것이 정답이므로 빈 문자열로 대체해도 결과가 같고,
                    # 매 검색마다 캐시 본문 전체(수백MB~수GB)를 메모리로 올리는 비용이 사라진다.
                    meta = self.text_cache.load_meta_for_prefix(self.path)
                    matched = self.text_cache.load_matching_texts(self.path, self.keyword)
                    self.file_cache = {
                        path: (mtime, size, matched.get(path, ""))
                        for path, (mtime, size) in meta.items()
                    }
            self.timing['cache_load'] = time.time() - t_phase_start

            # 실제로 더 이상 존재하지 않는 파일의 캐시 항목 정리
            # (zip경로!내부파일명 형태는 zip 파일 자체의 존재 여부로 판단)
            t_phase_start = time.time()
            exists_checked = {}
            stale_keys = []
            for cache_key in self.file_cache:
                outer_path = cache_key.split('!', 1)[0] if '!' in cache_key else cache_key
                if outer_path not in exists_checked:
                    exists_checked[outer_path] = Path(outer_path).exists()
                if not exists_checked[outer_path]:
                    stale_keys.append(cache_key)
            if stale_keys:
                for key in stale_keys:
                    del self.file_cache[key]
                self.text_cache.delete_entries(stale_keys)
            self.timing['cache_cleanup'] = time.time() - t_phase_start

            # CPU 바운드 텍스트 추출을 프로세스로 미리 처리
            # (GIL 하 스레드는 파싱을 사실상 직렬화함: 실측 결과 16스레드로도 1.03배에 그침)
            # 캐시에 이미 있는 항목은 대상에서 제외하고, 추출 결과는 file_cache에 채워 넣어
            # 이어지는 검색 단계(스레드)가 재파싱 없이 그대로 재사용하게 한다.
            # .doc/.ppt는 내용 추출 자체가 불가하므로 제외한다.
            t_phase_start = time.time()
            if not self.name_only:
                extract_targets = []
                extract_bytes = 0
                zip_files = []
                for f in files:
                    ext = f.suffix.lower()
                    if ext == '.zip':
                        zip_files.append(f)
                        continue
                    if ext in ('.doc', '.ppt'):
                        continue
                    try:
                        stat = f.stat()
                    except OSError:
                        continue
                    if self.skip_large_files and stat.st_size > 10 * 1024 * 1024:
                        continue
                    cache_key = (stat.st_mtime, stat.st_size)
                    cached = self.file_cache.get(str(f))
                    if cached and (cached[0], cached[1]) == cache_key:
                        continue
                    extract_targets.append(f)
                    extract_bytes += stat.st_size

                zip_tasks, zip_entry_count = self._collect_zip_extract_tasks(zip_files)
                if self.stop_flag:
                    self.status.emit("⏹️ 검색 중단됨")
                    self.finished.emit([])
                    return

                # 자식 프로세스가 PyQt5와 파서 라이브러리를 모두 import하므로 스폰 비용이 약 0.6초다.
                # 대상이 적으면 이 비용이 파싱 시간보다 커서 손해이므로(실측: 4개 18MB 기준
                # 스레드 0.53초 vs 프로세스풀 1.27초, 8개 31MB부터 프로세스풀이 유리) 그때는
                # 사전 추출을 건너뛰고 이어지는 검색 단계의 스레드가 처리하게 둔다.
                total_targets = len(extract_targets) + zip_entry_count
                worth_pool = total_targets >= 8 or extract_bytes >= 32 * 1024 * 1024

                # 추출량이 많으면 색인 갱신을 미뤘다가 마지막에 한 번에 만든다.
                # 이 조건(새 항목이 1,000건 이상이고 기존 캐시 항목 수 이상)은 사실상
                # "캐시DB가 비어 있는 상태에서 최초 1회 대량 검색"과 같은 상황이므로,
                # 캐시DB가 완성될 때까지 남은 시간을 상태표시바에 붉은 깜박임으로 알린다.
                defer_index = self.text_cache.should_defer_index(total_targets)
                if defer_index:
                    self.text_cache.begin_deferred_index()
                    self.extraction_urgent.emit(True)

                if total_targets and worth_pool:
                    self.status.emit(f"⚙️ {total_targets}개 파일 텍스트 추출 중...")
                    extract_start = time.monotonic()
                    extracted_count = 0
                    extracted_chars = 0
                    last_eta_emit = 0.0
                    with ProcessPoolExecutor(max_workers=self.max_workers) as pool:
                        extract_futures = {}
                        for f in extract_targets:
                            extract_futures[pool.submit(_extract_text_worker, str(f))] = False
                        for task in zip_tasks:
                            extract_futures[pool.submit(_extract_zip_entries_worker, task)] = True

                        for future in as_completed(extract_futures):
                            if self.stop_flag:
                                pool.shutdown(wait=False, cancel_futures=True)
                                if defer_index:
                                    self.extraction_urgent.emit(False)
                                self.status.emit("⏹️ 검색 중단됨")
                                self._log_extract_failures()
                                self.finished.emit([])
                                return
                            try:
                                result = future.result()
                            except Exception:
                                continue
                            # zip 워커는 여러 항목을 묶어서, 일반 워커는 한 항목을 돌려준다
                            entries = result if extract_futures[future] else [result]
                            chars_this_batch = 0
                            for path_str, key0, key1, text in entries:
                                if key0 is None:
                                    continue
                                if text is None:
                                    self.extract_failed_files.append(path_str)
                                    text = ""
                                entry = (key0, key1, text)
                                self.file_cache[path_str] = entry
                                self.new_cache_entries[path_str] = entry
                                chars_this_batch += len(text)
                            if len(self.new_cache_entries) >= self.CACHE_FLUSH_SIZE:
                                self._flush_cache()

                            if defer_index:
                                extracted_count += len(entries)
                                extracted_chars += chars_this_batch
                                now = time.monotonic()
                                if now - last_eta_emit >= 1.0:
                                    last_eta_emit = now
                                    # 처리 초반 소수 건은 프로세스풀 워커 생성 지연(약 0.6초)이
                                    # 그대로 실려 순간 처리율이 왜곡된다. 최소 20건은 처리된
                                    # 뒤부터 추정해 첫 표시부터 크게 틀리지 않게 한다.
                                    elapsed = now - extract_start
                                    if extracted_count >= 20:
                                        extraction_remaining = (elapsed / extracted_count) * (total_targets - extracted_count)
                                        # 이 검색이 "완성"되려면 추출 뒤 색인 생성(rebuild)까지
                                        # 필요하므로, 지금까지의 평균 본문 크기로 전체 본문량을
                                        # 추정해 rebuild 예상 시간(실측 0.110초/M자)까지 더한다.
                                        avg_chars = extracted_chars / extracted_count
                                        estimated_total_chars = avg_chars * total_targets
                                        index_estimate = (estimated_total_chars / 1_000_000
                                                          * self.INDEX_BUILD_SEC_PER_MILLION_CHARS)
                                        eta_text = _format_remaining(extraction_remaining + index_estimate)
                                    else:
                                        eta_text = "계산 중..."
                                    self.status.emit(
                                        f"⏳ 캐시DB 생성 중: {extracted_count:,}/{total_targets:,}개 "
                                        f"(남은 시간 약 {eta_text})"
                                    )

                if defer_index:
                    # 남은 항목까지 저장한 뒤 색인을 한 번에 만든다.
                    # 이 단계도 캐시DB가 아직 "완성"되지 않은 상태이므로 깜박임은 유지한다.
                    self._flush_cache()
                    if extracted_count > 0:
                        # 실제로 저장된 전체 행 수(기존 캐시 + 이번 추출분) 기준으로
                        # rebuild가 훑을 본문량을 다시 추정한다(추출 중 추정보다 정확함).
                        avg_chars = extracted_chars / extracted_count
                        index_estimate = (avg_chars * self.text_cache.row_count() / 1_000_000
                                         * self.INDEX_BUILD_SEC_PER_MILLION_CHARS)
                    else:
                        index_estimate = 0.0
                    self._run_with_countdown(
                        self.text_cache.end_deferred_index, index_estimate, "🔧 검색 색인을 만드는 중"
                    )
                    self.extraction_urgent.emit(False)
                    if self.stop_flag:
                        self.status.emit("⏹️ 검색 중단됨")
                        self._log_extract_failures()
                        self.finished.emit([])
                        return
            self.timing['extract'] = time.time() - t_phase_start

            # 병렬 검색
            t_phase_start = time.time()
            results = []
            processed = 0
            checkpoint_logged = set()

            last_pct = -1

            # 파일마다 태스크를 만들면 future 7천여 개의 관리 비용만 0.15초가 든다.
            # 스레드는 zip 열기 같은 I/O 병렬화에 여전히 필요하므로 풀은 유지하고 묶어서 제출한다.
            chunk_size = max(1, min(64, len(files) // (self.max_workers * 4)))
            chunks = [files[i:i + chunk_size] for i in range(0, len(files), chunk_size)]

            def search_chunk(chunk):
                chunk_results = []
                for f in chunk:
                    if self.stop_flag:
                        break
                    chunk_results.extend(self._search_file(f, regex))
                return chunk_results

            with ThreadPoolExecutor(max_workers=self.max_workers) as executor:
                futures = {executor.submit(search_chunk, c): len(c) for c in chunks}
                stop_requested = False

                for future in as_completed(futures):
                    try:
                        file_results = future.result()
                        results.extend(file_results)
                        # 실시간으로 각 결과 전송
                        for result in file_results:
                            self.result_found.emit(result)
                    except Exception as e:
                        logging.getLogger(__name__).warning(f"파일 청크 검색 실패: {e}")

                    # 중단 플래그 확인 (사용자가 중지 버튼을 눌렀거나, _report_found()가
                    # 결과 MAX_RESULTS 초과를 감지해 자동으로 세운 경우 둘 다 여기서 걸린다).
                    # 여기서 바로 return하지 않는 이유: 다른 스레드가 방금 stop_flag를
                    # 세운 시점에 그 스레드 자신의 future(예: zip 내부를 순회하다 막
                    # 멈춘 작업)는 아직 as_completed에 나타나지 않았을 수 있다. 즉시
                    # 종료하면 그 future의 결과가 통째로 유실된다(실측으로 확인한 버그).
                    # cancel_futures=True로 아직 시작 안 한 작업만 취소하고, 이미 실행
                    # 중이던 작업들은 stop_flag 덕분에 금방 끝나므로 루프를 끝까지 돌려
                    # 그 결과까지 마저 수거한 뒤 종료한다.
                    if self.stop_flag and not stop_requested:
                        stop_requested = True
                        executor.shutdown(wait=False, cancel_futures=True)

                    # zip 내부 파일 추출 결과가 쌓이므로 여기서도 주기적으로 저장한다
                    if len(self.new_cache_entries) >= self.CACHE_FLUSH_SIZE:
                        self._flush_cache()

                    processed += futures[future]
                    pct = int((processed / len(files)) * 100)
                    # 진행바는 정수 퍼센트 단위라 값이 바뀔 때만 보내면 표시 결과가 같다
                    # (파일마다 보내면 7천여 회 크로스스레드 시그널이 발생한다)
                    if pct != last_pct:
                        last_pct = pct
                        self.progress.emit(pct)

                    # 진단용: 25/50/75/100% 지점에서 경과시간 기록
                    for cp in (25, 50, 75, 100):
                        if pct >= cp and cp not in checkpoint_logged:
                            checkpoint_logged.add(cp)
                            self.timing[f'checkpoint_{cp}pct'] = time.time() - t_phase_start

                if stop_requested:
                    if self._max_results_hit:
                        self.status.emit(
                            f"⚠️ 결과가 {self.MAX_RESULTS}건을 초과해 검색을 중단했습니다. "
                            "검색어를 더 구체적으로 입력해주세요."
                        )
                    else:
                        self.status.emit("⏹️ 검색 중단됨")
                    self._log_extract_failures()
                    self.finished.emit(results)
                    return
            self.timing['search'] = time.time() - t_phase_start

            # 매칭된 파일들과 미매칭 파일들 분류
            t_phase_start = time.time()
            matched_files = {result['full_path'] for result in results}
            no_match_files = [str(f) for f in files if str(f) not in matched_files]
            self.timing['classify'] = time.time() - t_phase_start

            self.timing['total'] = time.time() - t_run_start
            self._log_extract_failures()
            self.finished.emit(results)
            self.no_match_files_updated.emit(no_match_files)
            self.skipped_files_updated.emit(self.skipped_files_list)
            self.status.emit(f"✅ 검색 완료: {len(results)}개 결과")

        except Exception as e:
            self.error.emit(f"오류 발생: {str(e)}")
        finally:
            self._flush_cache()
            # 대량 추출 도중 예외가 나도 상태표시바 깜박임이 계속 남지 않게 한다
            # (정상 종료 시 이미 꺼져 있으므로 다시 꺼도 무해하다)
            self.extraction_urgent.emit(False)

    def _collect_files(self) -> List[Path]:
        """파일 수집 - 지정된 파일 형식만"""
        # 허용된 파일 확장자 (zip은 내부 파일을 풀어서 검색)
        allowed_extensions = SearchWorker.EXTRACTABLE_EXTENSIONS | {'.zip'}

        files = []
        self.excluded_files = []  # 초기화

        # ignore_dirs를 정규화된 Path 객체로 변환
        excluded_paths = {Path(excluded).resolve() for excluded in self.ignore_dirs}

        for root, dirs, filenames in os.walk(self.path):
            # 중단 플래그 확인
            if self.stop_flag:
                return files

            # 제외할 폴더 필터링 - 전체 경로로 비교
            # resolve()는 디렉터리마다 발생하는 syscall이라 비용이 크므로(3,294개 기준 0.76초)
            # 제외 폴더가 실제로 설정된 경우에만 수행한다
            if excluded_paths:
                dirs_to_remove = []
                for d in dirs:
                    dir_path = Path(root) / d
                    dir_path_resolved = dir_path.resolve()

                    # 제외 폴더 또는 그 하위 폴더인지 확인
                    for excluded_path in excluded_paths:
                        try:
                            # dir_path가 excluded_path의 하위인지 확인
                            dir_path_resolved.relative_to(excluded_path)
                            dirs_to_remove.append(d)
                            break
                        except ValueError:
                            # 하위가 아니면 continue
                            pass

                # 제외할 폴더 제거
                for d in dirs_to_remove:
                    if d in dirs:
                        dirs.remove(d)

            # 경로 정규화는 파일마다가 아니라 디렉터리마다 한 번만 수행한다
            root_path = Path(root)
            root_display = str(root_path)

            for filename in filenames:
                # 확장자 판정은 문자열로 처리하고 대상 파일만 Path 객체를 만든다
                # (전체 6만여 개 파일에 Path를 만들면 0.39초가 더 든다)
                ext = SearchWorker._entry_ext(filename)
                if ext in allowed_extensions:
                    files.append(root_path / filename)
                else:
                    # 제외된 파일 추적
                    self.excluded_files.append(os.path.join(root_display, filename))

        # 제외된 파일 목록 업데이트 신호
        self.excluded_files_updated.emit(self.excluded_files)

        return files

    # 검색이 중간에 끊겨도 진행분이 남도록 이 개수마다 캐시를 디스크에 저장한다.
    # 저장 자체에 비용이 있으므로(FTS 인덱스 갱신 포함) 너무 잦지 않게 잡는다.
    CACHE_FLUSH_SIZE = 500

    # FTS5 rebuild 소요 시간 실측치(본문 1M자당 초). 캐시DB가 비어 있어 색인 생성을
    # 뒤로 미룬 대량 추출 중 "캐시DB 완성까지" 남은 시간을 추정하는 데 쓴다.
    INDEX_BUILD_SEC_PER_MILLION_CHARS = 0.110

    def _run_with_countdown(self, fn, estimate_seconds, label):
        """진행률을 알 수 없는 단일 작업(fn, 예: FTS rebuild)을 실행하는 동안,
        미리 추정한 소요 시간에서 경과 시간만큼 줄여가며 상태표시바에 남은 시간을 보여준다.
        fn은 이 메서드를 호출한 스레드에서 그대로(동기) 실행되고, 별도 스레드는 오직
        1초 간격으로 상태 텍스트만 emit한다. fn에는 stop_flag를 확인하는 콜백을 넘겨,
        도중에 중지 버튼이 눌리면 fn 내부에서 SQL 실행 자체를 중단시킬 수 있게 한다."""
        stop_event = threading.Event()

        def ticker():
            start = time.monotonic()
            while not stop_event.wait(1.0):
                remaining = max(0.0, estimate_seconds - (time.monotonic() - start))
                self.status.emit(f"{label}... (남은 시간 약 {_format_remaining(remaining)})")

        self.status.emit(f"{label}... (남은 시간 약 {_format_remaining(estimate_seconds)})")
        ticker_thread = threading.Thread(target=ticker, daemon=True)
        ticker_thread.start()
        try:
            fn(stop_check=lambda: self.stop_flag)
        finally:
            stop_event.set()
            ticker_thread.join(timeout=2.0)

    # zip 하나에 엔트리가 수천 개인 경우가 있어, 한 작업이 너무 커지지 않도록 나눈다
    # (워커가 돌려주는 텍스트가 한꺼번에 메모리에 올라오고 부하 분배도 나빠진다).
    # 너무 잘게 나누면 작업마다 zip을 다시 열어 중앙 디렉터리를 읽는 비용이 붙는다.
    ZIP_TASK_CHUNK = 100

    def _collect_zip_extract_tasks(self, zip_files):
        """zip 내부 파일 중 캐시에 없는 것들을 프로세스풀 작업 단위로 묶는다.
        반환: ([(zip경로, [(내부파일명, 확장자), ...]), ...], 총 엔트리 수)"""
        tasks = []
        total = 0
        for zip_path in zip_files:
            if self.stop_flag:
                break
            zip_str = str(zip_path)
            try:
                zip_stat = zip_path.stat()
            except OSError:
                continue
            zip_key = (zip_stat.st_mtime, zip_stat.st_size)

            # zip 자체가 그대로이고 내부 파일이 모두 캐시된 적이 있으면 중앙 디렉터리를
            # 다시 읽지 않는다. 이 열거는 zip 298개 기준 0.4초로, 캐시가 다 찬 뒤에도
            # 매 검색마다 발생하면 그만큼 손해다.
            marker = self.file_cache.get(zip_str)
            if marker and (marker[0], marker[1]) == zip_key:
                continue

            try:
                with zipfile.ZipFile(zip_path) as zf:
                    infos = zf.infolist()
            except Exception as e:
                logging.getLogger(__name__).warning(f"zip 파일 열기 실패({zip_str}): {e}")
                continue

            pending = []
            for info in infos:
                if info.is_dir():
                    continue
                ext = self._entry_ext(info.filename)
                # .doc/.ppt는 내용 추출을 지원하지 않아 검색 단계에서도 본문을 보지 않는다
                if ext not in SearchWorker.EXTRACTABLE_EXTENSIONS or ext in ('.doc', '.ppt'):
                    continue
                if self.skip_large_files and info.file_size > 10 * 1024 * 1024:
                    continue
                # 캐시 키는 검색 단계(_extract_zip_entry_text)와 동일한 형식이어야 한다
                cached = self.file_cache.get(f"{zip_str}!{info.filename}")
                if cached and (cached[0], cached[1]) == (info.CRC, info.file_size):
                    continue
                pending.append((info.filename, ext))

            if pending:
                for i in range(0, len(pending), self.ZIP_TASK_CHUNK):
                    tasks.append((zip_str, pending[i:i + self.ZIP_TASK_CHUNK]))
                total += len(pending)
            else:
                # 내부 파일이 모두 캐시돼 있으므로 다음 검색부터는 열지 않아도 된다는 표시.
                # zip 자체는 텍스트 추출 대상이 아니라 이 키가 본문으로 읽히는 일은 없다.
                entry = (zip_key[0], zip_key[1], "")
                self.file_cache[zip_str] = entry
                self.new_cache_entries[zip_str] = entry

        return tasks, total

    def _flush_cache(self):
        """지금까지 추출한 항목을 디스크에 저장하고 버퍼를 비운다.
        저장은 검색이 끝난 뒤 한 번만 하면 중단 시 그동안의 추출 결과가 전부 사라지므로
        (D:/ 전체처럼 오래 걸리는 검색에서는 매번 처음부터 다시 추출하게 됨) 중간중간 호출한다."""
        with self.cache_lock:
            if not self.new_cache_entries:
                return
            pending = self.new_cache_entries
            self.new_cache_entries = {}
        # 디스크 저장은 느리므로 락을 놓은 뒤 수행한다
        self.text_cache.save_entries(pending)

    def _extract_text(self, file_path: Path) -> str:
        """파일 형식별로 텍스트 추출 (디스크 캐시 지원 - 재검색 시 재사용)"""
        file_path_str = str(file_path)

        try:
            stat = file_path.stat()
            cache_key = (stat.st_mtime, stat.st_size)
        except OSError:
            return _extract_text_impl(file_path)

        # 캐시 확인 (파일이 변경되지 않았으면 재사용)
        cached = self.file_cache.get(file_path_str)
        if cached and (cached[0], cached[1]) == cache_key:
            return cached[2]

        text = _extract_text_impl(file_path)
        if text is None:
            self.extract_failed_files.append(file_path_str)
            text = ""

        with self.cache_lock:
            entry = (cache_key[0], cache_key[1], text)
            self.file_cache[file_path_str] = entry
            self.new_cache_entries[file_path_str] = entry

        return text

    def _report_found(self, n=1):
        """결과 n건이 방금 발견됐음을 전역 카운터에 반영하고, MAX_RESULTS를 넘으면
        stop_flag를 세운다. zip 하나에 매치가 수천 건 몰려 있으면 그 zip을 다 처리할
        때까지는 메인 루프의 len(results) 체크가 전혀 반영되지 않으므로(zip 처리가
        끝나야 결과가 한꺼번에 합쳐짐), 매칭이 발견되는 즉시(엔트리 단위로) 여기서
        전역으로 집계해야 zip 내부 순회 중에도 바로 멈출 수 있다."""
        with self._found_lock:
            self._found_count += n
            if self._found_count >= self.MAX_RESULTS:
                self._max_results_hit = True
                self.stop_flag = True

    def _log_extract_failures(self):
        """텍스트 추출(파싱) 자체가 실패한 파일이 있으면 로그로 남긴다. 검색 결과에는
        영향 없이 빈 문자열로 취급해 조용히 넘어가지만, 원인 추적을 위해 기록해둔다."""
        if self.extract_failed_files:
            logging.getLogger(__name__).warning(
                f"텍스트 추출 실패 {len(self.extract_failed_files)}건: "
                + ', '.join(self.extract_failed_files)
            )

    def _search_file(self, file_path: Path, regex):
        """단일 파일 검색"""
        results = []

        # 현재 처리 중인 파일 신호 전송
        # 파일마다 보내면 초당 수천 건이라 화면에서는 어차피 읽을 수 없고 시그널 비용만 커지므로
        # 표시 간격을 최소 30ms로 둔다 (스레드 간 경합이 나도 발행이 몇 번 늘 뿐 문제되지 않음)
        now = time.monotonic()
        if now - self._last_file_signal >= 0.03:
            self._last_file_signal = now
            self.file_processing.emit(f"🔍 {file_path.name}")

        # 파일 정보 수집
        try:
            stat = file_path.stat()
            file_size = stat.st_size
            mod_time = time.strftime('%Y-%m-%d %H:%M:%S', time.localtime(stat.st_mtime))
        except:
            file_size = 0
            mod_time = "Unknown"

        # 대용량 파일 스킵 (10MB 이상)
        if self.skip_large_files and file_size > 10 * 1024 * 1024:
            self.skipped_large_files += 1
            self.skipped_files_list.append(str(file_path))
            self.skipped_files_count.emit(self.skipped_large_files)
            return results

        # zip 압축파일은 내부 파일을 풀어서 검색
        if file_path.suffix.lower() == '.zip':
            return self._search_zip_file(file_path, regex, file_size, mod_time)

        # 파일명과 폴더 경로 분리
        filename = file_path.name
        folder_path = str(file_path.parent)

        # 파일 아이콘 가져오기
        icon = SearchWorker.get_file_icon(filename)
        filename_with_icon = f"{icon} {filename}"

        total_match_count = 0
        filename_matched = False
        matched_lines = []  # 매칭된 라인 저장

        # .doc/.ppt(구버전 워드/파워포인트)는 내용 추출을 지원하지 않으므로 검색 옵션과 무관하게 항상 파일명만 매칭
        content_unsupported = file_path.suffix.lower() in ('.doc', '.ppt')

        # 파일명 검색
        if not self.content_only or content_unsupported:
            if self._match_keyword(file_path.name, regex):
                total_match_count += 1
                filename_matched = True

        # 파일 내용 검색 - 파일 형식별로 텍스트 추출
        content_match_count = 0
        if not self.name_only and not content_unsupported:
            # 파일에서 텍스트 추출
            text = self._extract_text(file_path)
            if text:
                if regex:
                    for line in text.split('\n'):
                        if regex.search(line):
                            total_match_count += 1
                            content_match_count += 1
                            # 매칭된 라인 저장 (최대 4줄)
                            if len(matched_lines) < 4:
                                matched_lines.append(line.strip())
                else:
                    # OR-그룹 중 하나라도 AND-항이 전부 있는지 문서 전체 기준으로 먼저 확인
                    if _query_matches(self.keyword_query, text.lower()):
                        for line in text.split('\n'):
                            # 라인 단위로는 AND 전체 충족이 아니라 쿼리 항 중 하나라도 있으면
                            # 미리보기/카운트 대상으로 삼는다(AND 항들이 서로 다른 줄에 있을 수 있음)
                            if _query_any_term_in(self.keyword_query, line.lower()):
                                total_match_count += 1
                                content_match_count += 1
                                # 매칭된 라인 저장 (최대 4줄)
                                if len(matched_lines) < 4:
                                    matched_lines.append(line.strip())

        # 매칭이 있으면 파일당 하나의 결과만 추가
        if total_match_count > 0:
            results.append({
                'type': 'file_summary',
                'filename': filename_with_icon,
                'folder_path': folder_path,
                'full_path': str(file_path),
                'size': file_size,
                'modified': mod_time,
                'match_count': total_match_count,
                'matched_lines': matched_lines,  # 매칭된 라인 추가
                'extension': file_path.suffix.lower()  # 파일 확장자 추가
            })
            self._report_found()

        return results

    @staticmethod
    def _fix_zip_entry_name(info) -> str:
        """UTF-8 플래그 없이 CP949로 저장된 zip 내부 파일명을 복원 (한글 zip 도구에서 흔한 인코딩 문제)"""
        if info.flag_bits & 0x800:  # UTF-8 플래그가 설정된 경우 zipfile이 이미 올바르게 디코딩함
            return info.filename
        try:
            return info.filename.encode('cp437').decode('cp949')
        except (UnicodeDecodeError, UnicodeEncodeError):
            return info.filename

    @staticmethod
    def _entry_basename(entry_name: str) -> str:
        """zip 엔트리 경로에서 파일명만 분리 (Path.name과 동일, Path 객체 생성 없이)"""
        name = entry_name.rstrip('/\\')
        sep = max(name.rfind('/'), name.rfind('\\'))
        return name[sep + 1:] if sep >= 0 else name

    @staticmethod
    def _entry_ext(entry_name: str) -> str:
        """zip 엔트리의 소문자 확장자 (Path.suffix와 동일한 규칙:
        점으로 시작하는 파일명과 점으로 끝나는 파일명은 확장자가 없는 것으로 본다)"""
        name = SearchWorker._entry_basename(entry_name)
        dot = name.rfind('.')
        return name[dot:].lower() if 0 < dot < len(name) - 1 else ''

    def _search_zip_file(self, file_path: Path, regex, zip_size, mod_time):
        """zip 압축파일 내부 문서를 풀어서 검색 (내부 파일마다 결과 하나씩)"""
        results = []
        folder_path = str(file_path.parent)
        zip_icon = SearchWorker.get_file_icon(file_path.name)

        try:
            with zipfile.ZipFile(file_path) as zf:
                for info in zf.infolist():
                    if self.stop_flag:
                        break
                    if info.is_dir():
                        continue
                    # 확장자는 ASCII라 CP437->CP949 보정 전후가 같으므로, 비용이 큰 이름 보정과
                    # Path 객체 생성보다 먼저 확장자로 걸러낸다
                    # (엔트리 9만여 개 중 대상은 3만여 개, 이 순서만 바꿔도 0.35초가 준다)
                    ext = self._entry_ext(info.filename)
                    if ext not in SearchWorker.EXTRACTABLE_EXTENSIONS:
                        continue
                    if self.skip_large_files and info.file_size > 10 * 1024 * 1024:
                        continue
                    # UTF-8 플래그 없이 CP949로 저장된 한글 zip 파일명 복원(한글 zip 도구에서 흔함)
                    fixed_filename = self._fix_zip_entry_name(info)
                    inner_name = self._entry_basename(fixed_filename)

                    match_count = 0
                    matched_lines = []

                    # .doc/.ppt(구버전 워드/파워포인트)는 내용 추출을 지원하지 않으므로 검색 옵션과 무관하게 항상 파일명만 매칭
                    content_unsupported = ext in ('.doc', '.ppt')

                    # 내부 파일명 검색
                    if not self.content_only or content_unsupported:
                        if self._match_keyword(inner_name, regex):
                            match_count += 1

                    # 내부 파일 내용 검색
                    if not self.name_only and not content_unsupported:
                        text = self._extract_zip_entry_text(file_path, zf, info, ext)
                        if text:
                            if regex:
                                for line in text.split('\n'):
                                    if regex.search(line):
                                        match_count += 1
                                        if len(matched_lines) < 4:
                                            matched_lines.append(line.strip())
                            else:
                                if _query_matches(self.keyword_query, text.lower()):
                                    for line in text.split('\n'):
                                        if _query_any_term_in(self.keyword_query, line.lower()):
                                            match_count += 1
                                            if len(matched_lines) < 4:
                                                matched_lines.append(line.strip())

                    if match_count > 0:
                        inner_icon = SearchWorker.get_file_icon(inner_name)
                        results.append({
                            'type': 'file_summary',
                            'filename': f"{zip_icon} {file_path.name} → {inner_icon} {inner_name}",
                            'folder_path': folder_path,
                            'full_path': str(file_path),
                            'size': zip_size,
                            'modified': mod_time,
                            'match_count': match_count,
                            'matched_lines': matched_lines,
                            'extension': ext,
                            'inner_name': inner_name
                        })
                        self._report_found()
        except Exception as e:
            logging.getLogger(__name__).warning(f"zip 파일 검색 실패({file_path}): {e}")

        return results

    def _extract_zip_entry_text(self, zip_path: Path, zf, info, ext) -> str:
        """zip 내부 파일의 텍스트 추출 (메모리에서 바로 파싱, 캐싱 지원)"""
        cache_key_str = f"{zip_path}!{info.filename}"
        cache_key = (info.CRC, info.file_size)

        cached = self.file_cache.get(cache_key_str)
        if cached and (cached[0], cached[1]) == cache_key:
            return cached[2]

        try:
            text = _extract_text_from_bytes(zf.read(info.filename), ext)
        except Exception:
            text = None

        if text is None:
            self.extract_failed_files.append(cache_key_str)
            text = ""

        with self.cache_lock:
            entry = (cache_key[0], cache_key[1], text)
            self.file_cache[cache_key_str] = entry
            self.new_cache_entries[cache_key_str] = entry

        return text

    def _match_keyword(self, text: str, regex):
        """키워드 매칭"""
        if regex:
            return regex.search(text) is not None
        else:
            return _query_matches(self.keyword_query, text.lower())

    def _is_binary(self, path: Path) -> bool:
        """바이너리 파일 확인"""
        binary_exts = {'exe', 'dll', 'so', 'png', 'jpg', 'zip', 'db', 'pdf'}
        return path.suffix.lower().lstrip('.') in binary_exts

    @staticmethod
    def get_file_icon(filename: str) -> str:
        """파일 형식에 따라 아이콘 반환"""
        ext = Path(filename).suffix.lower()

        icon_map = {
            # 문서
            '.txt': '📄',
            '.md': '📝',
            '.pdf': '📕',
            '.doc': '📘',
            '.docx': '📘',
            '.ppt': '📙',
            '.pptx': '📙',
            '.hwp': '📗',
            '.hwpx': '📗',
            '.xls': '📊',
            '.xlsx': '📊',
            '.csv': '📋',
            '.json': '⚙️',
            '.xml': '⚙️',
            '.yaml': '⚙️',
            '.yml': '⚙️',

            # 코드
            '.py': '🐍',
            '.js': '⚡',
            '.ts': '📘',
            '.tsx': '⚛️',
            '.jsx': '⚛️',
            '.cpp': '⚙️',
            '.c': '⚙️',
            '.h': '⚙️',
            '.java': '☕',
            '.cs': '💠',
            '.go': '🐹',
            '.rs': '🦀',
            '.rb': '💎',
            '.php': '🐘',
            '.html': '🌐',
            '.css': '🎨',
            '.scss': '🎨',
            '.sql': '🗄️',
            '.sh': '🔧',
            '.bat': '🔧',
            '.ps1': '🔧',

            # 이미지
            '.png': '🖼️',
            '.jpg': '🖼️',
            '.jpeg': '🖼️',
            '.gif': '🎬',
            '.svg': '🎨',
            '.ico': '🎯',

            # 압축
            '.zip': '📦',
            '.rar': '📦',
            '.7z': '📦',
            '.tar': '📦',
            '.gz': '📦',

            # 실행파일
            '.exe': '⚙️',
            '.dll': '⚙️',
            '.so': '⚙️',

            # 기타
            '.git': '🔀',
            '.env': '🔐',
            '.log': '📋',
        }

        return icon_map.get(ext, '📁')


class FSearchGUI(QMainWindow):
    """fsearch GUI 메인 윈도우"""

    def __init__(self):
        super().__init__()
        self.search_worker = None
        self.is_searching = False  # 검색 진행 중 플래그
        self.results = []
        self.sort_column = None  # 헤더 클릭 정렬 기준 컬럼 (None이면 업무 연관도순)
        self.sort_ascending = True
        self.excluded_files = []  # 제외된 파일 목록
        self.read_files = []  # 읽은 파일 목록 (누적)
        self.no_match_files = []  # 검색어 미포함 파일 목록 (누적)
        self.skipped_files = []  # 스킵된 대용량파일 목록 (누적)
        self.logger = setup_logging()  # 로깅 설정
        self.search_history = SearchHistory()  # 검색 이력 관리
        self.search_start_time = None  # 검색 시작 시간
        self.search_elapsed_time = 0  # 검색 소요 시간 (초)
        self.skipped_files_count_total = 0  # 스킵된 파일 갯수
        self._icon_provider = QFileIconProvider()  # 윈도우 탐색기 파일 아이콘 조회
        self._icon_html_cache = {}  # 확장자별 아이콘 <img> HTML 캐시
        self._prewarm_icon_cache()  # 검색 중 첫 조회 지연(확장자당 최대 수백ms)을 시작 시점으로 이동
        self.blink_timer = QTimer()  # 버튼 깜박임 타이머
        self.blink_timer.timeout.connect(self.toggle_button_blink)  # 타이머 신호 연결
        self.blink_state = False  # 깜박임 상태

        self.status_blink_timer = QTimer()  # 상태표시바 깜박임 타이머(캐시DB 생성 중 표시용)
        self.status_blink_timer.timeout.connect(self.toggle_status_blink)
        self.status_blink_state = False
        self.init_ui()

    def init_ui(self):
        """UI 초기화"""
        self.setWindowTitle("🔍 fsearch - 파일 검색 도구")
        self.setGeometry(100, 100, 1200, 600)

        # 중앙 위젯
        central = QWidget()
        self.setCentralWidget(central)
        layout = QVBoxLayout(central)
        layout.setContentsMargins(5, 5, 5, 5)  # 마진 최소화
        layout.setSpacing(3)  # 요소 간 간격 축소

        # ===== 검색 옵션 영역 =====
        options_layout = QHBoxLayout()

        # 경로 선택 (경로 이력 드롭다운)
        options_layout.addWidget(QLabel("경로:"))
        self.path_input = QComboBox()
        self.path_input.setEditable(True)
        self.path_input.setMaximumHeight(25)

        # 저장된 경로 로드 (드롭다운 목록만 채우고, 표시되는 기본값은 비워둠)
        saved_paths = self.search_history.get_paths()
        if saved_paths:
            self.path_input.addItems(saved_paths)
        self.path_input.setCurrentText("")

        options_layout.addWidget(self.path_input)

        browse_btn = QPushButton("찾아보기")
        browse_btn.clicked.connect(self.browse_path)
        browse_btn.setMaximumHeight(25)
        options_layout.addWidget(browse_btn)

        self.clear_path_history_btn = QPushButton("🗑️ 경로 초기화")
        self.clear_path_history_btn.clicked.connect(self.clear_path_history)
        self.clear_path_history_btn.setMaximumHeight(25)
        options_layout.addWidget(self.clear_path_history_btn)

        # 검색어 (검색 이력 드롭다운)
        options_layout.addWidget(QLabel("검색:"))
        self.keyword_input = QComboBox()
        self.keyword_input.setEditable(True)
        self.keyword_input.lineEdit().setPlaceholderText("검색할 키워드 입력...")
        self.keyword_input.lineEdit().setToolTip(
            "공백으로 구분하면 AND 조건(예: 네이버 회의록)\n"
            "OR을 명시하면 OR 조건(예: 네이버 OR 카카오)"
        )
        self.keyword_input.lineEdit().returnPressed.connect(self.search)
        self.keyword_input.lineEdit().textChanged.connect(self.on_keyword_changed)
        self.keyword_input.setMaximumHeight(25)

        # 검색 이력 로드
        keywords = self.search_history.get_keywords()
        if keywords:
            self.keyword_input.addItems(keywords)

        options_layout.addWidget(self.keyword_input, 2)

        self.search_btn = QPushButton("🔍 검색")
        self.search_btn.clicked.connect(self.search)
        self.search_btn.setMaximumHeight(25)
        options_layout.addWidget(self.search_btn)

        self.clear_history_btn = QPushButton("🗑️ 검색어 초기화")
        self.clear_history_btn.clicked.connect(self.clear_search_history)
        self.clear_history_btn.setMaximumHeight(25)
        options_layout.addWidget(self.clear_history_btn)

        layout.addLayout(options_layout)

        # ===== 경로 제외 영역 + 자주 제외하는 폴더 통계 (같은 행) =====
        combined_container = QWidget()
        combined_layout = QHBoxLayout(combined_container)
        combined_layout.setContentsMargins(0, 1, 0, 1)
        combined_layout.setSpacing(0)

        # --- 경로 제외 영역 ---
        exclude_container = QWidget()
        self.exclude_layout = QHBoxLayout(exclude_container)
        self.exclude_layout.setContentsMargins(0, 2, 0, 2)
        self.exclude_layout.setSpacing(5)

        exclude_label = QLabel("제외 폴더:")
        self.exclude_layout.addWidget(exclude_label)

        add_exclude_btn = QPushButton("+ 폴더 추가")
        add_exclude_btn.clicked.connect(self.add_exclude_folder)
        add_exclude_btn.setMaximumHeight(25)
        add_exclude_btn.setMaximumWidth(100)
        self.exclude_layout.addWidget(add_exclude_btn)

        # 제외 폴더 체크박스 저장소
        self.exclude_checkboxes = {}

        exclude_container.setMaximumHeight(30)

        combined_layout.addWidget(exclude_container)

        # --- 자주 제외하는 폴더 통계 영역 ---
        stats_container = QWidget()
        stats_layout = QHBoxLayout(stats_container)
        stats_layout.setContentsMargins(0, 2, 0, 2)
        stats_layout.setSpacing(5)

        stats_label = QLabel("📊 자주 제외하는 폴더:")
        stats_layout.addWidget(stats_label)

        # 자주 제외하는 폴더 표시
        self.stats_label = QLabel()
        self.update_excluded_stats_display()
        stats_layout.addWidget(self.stats_label)

        stats_container.setMaximumHeight(30)

        combined_layout.addWidget(stats_container)

        # 남은 공간 채우기
        combined_layout.addStretch()

        combined_container.setMaximumHeight(30)
        layout.addWidget(combined_container)

        # ===== 추가 옵션 영역 (검색 설정 + 파일 통계 버튼) =====
        # 한 줄에 배치하되, 창이 좁아지면 FlowLayout이 자동으로 다음 줄로 줄바꿈함
        options2_container = QWidget()
        options2_layout = FlowLayout(options2_container)
        options2_layout.setContentsMargins(0, 0, 0, 0)

        self.name_only_cb = QCheckBox("파일명만")
        self.name_only_cb.setMaximumHeight(25)
        self.content_only_cb = QCheckBox("내용만")
        self.content_only_cb.setMaximumHeight(25)
        self.regex_cb = QCheckBox("정규식")
        self.regex_cb.setMaximumHeight(25)
        self.skip_large_cb = QCheckBox("대용량파일 스킵(>10MB)")
        self.skip_large_cb.setMaximumHeight(25)

        options2_layout.addWidget(self.name_only_cb)
        options2_layout.addWidget(self.content_only_cb)
        options2_layout.addWidget(self.regex_cb)
        options2_layout.addWidget(self.skip_large_cb)

        options2_layout.addWidget(QLabel("스레드:"))
        self.workers_spin = QSpinBox()
        self.workers_spin.setValue(16)
        self.workers_spin.setMinimum(1)
        self.workers_spin.setMaximum(32)
        self.workers_spin.setMaximumHeight(25)
        options2_layout.addWidget(self.workers_spin)

        self.excluded_btn = QPushButton("📁 제외된 파일")
        self.excluded_btn.clicked.connect(self.show_excluded_files)
        self.excluded_btn.setMaximumHeight(25)
        options2_layout.addWidget(self.excluded_btn)

        self.read_files_btn = QPushButton("🔍 찾은 파일 수 (0)")
        self.read_files_btn.clicked.connect(self.show_read_files)
        self.read_files_btn.setMaximumHeight(25)
        options2_layout.addWidget(self.read_files_btn)

        self.no_match_files_btn = QPushButton("❌ 검색어 미포함 파일 수 (0)")
        self.no_match_files_btn.clicked.connect(self.show_no_match_files)
        self.no_match_files_btn.setMaximumHeight(25)
        options2_layout.addWidget(self.no_match_files_btn)

        self.skipped_files_btn = QPushButton("⏭️ 스킵된 대용량파일 (0)")
        self.skipped_files_btn.clicked.connect(self.show_skipped_files)
        self.skipped_files_btn.setMaximumHeight(25)
        self.skipped_files_btn.setEnabled(False)
        options2_layout.addWidget(self.skipped_files_btn)

        self.performance_btn = QPushButton("⏱️ 완료시간 (0.00초)")
        self.performance_btn.setMaximumHeight(25)
        self.performance_btn.setEnabled(False)
        options2_layout.addWidget(self.performance_btn)

        self.export_excel_btn = QPushButton("📊 엑셀로 내보내기")
        self.export_excel_btn.clicked.connect(self.export_to_excel)
        self.export_excel_btn.setMaximumHeight(25)
        options2_layout.addWidget(self.export_excel_btn)

        layout.addWidget(options2_container)

        # ===== 진행바 =====
        self.progress_bar = QProgressBar()
        self.progress_bar.setMaximum(100)
        self.progress_bar.setVisible(False)
        self.progress_bar.setMaximumHeight(20)
        layout.addWidget(self.progress_bar)

        # ===== 상태 메시지 =====
        self.status_label = QLabel("준비 완료")
        self.status_label.setMaximumHeight(20)
        layout.addWidget(self.status_label)

        # ===== 결과 필터 (확장자/날짜) =====
        # self.results 자체는 건드리지 않고 표시(테이블/텍스트 탭)만 좁힌다.
        # 확장자 체크박스(검색할 때마다 새로 구성)와 날짜/초기화(고정)를 같은 FlowLayout
        # 한 줄에 순서대로 배치해, 마지막 체크박스 바로 옆에 날짜 필터가 붙어 보이게 한다.
        # 체크박스는 검색마다 없앴다 새로 만들지만 날짜/초기화 위젯은 그대로 두고 그때마다
        # 맨 끝으로 다시 옮겨 붙인다(_rebuild_extension_filter 참고).
        filter_container = QWidget()
        self.filter_ext_layout = FlowLayout(filter_container, spacing=5)
        self.filter_ext_layout.addWidget(QLabel("확장자:"))
        self.filter_ext_checkboxes = {}  # 확장자 -> QCheckBox (검색할 때마다 새로 구성)

        # 날짜를 한 글자씩 입력할 때마다 필터를 바로 적용하면(특히 결과가 많을 때) 매
        # 키 입력마다 다시 그리느라 입력이 밀리므로, 타이핑이 잠깐 멈췄을 때 한 번만
        # 적용되도록 디바운스한다.
        self._filter_debounce_timer = QTimer(self)
        self._filter_debounce_timer.setSingleShot(True)
        self._filter_debounce_timer.timeout.connect(self._apply_filters)

        self.filter_date_from = QLineEdit()
        self.filter_date_from.setPlaceholderText("YYYY-MM-DD")
        self.filter_date_from.setMaximumWidth(90)
        self.filter_date_from.textChanged.connect(self._schedule_apply_filters)
        self.filter_date_to = QLineEdit()
        self.filter_date_to.setPlaceholderText("YYYY-MM-DD")
        self.filter_date_to.setMaximumWidth(90)
        self.filter_date_to.textChanged.connect(self._schedule_apply_filters)
        filter_reset_btn = QPushButton("필터 초기화")
        filter_reset_btn.setMaximumHeight(25)
        filter_reset_btn.clicked.connect(self.reset_filters)
        # 체크박스 뒤에 이어 붙일 고정 위젯들(순서대로) - 매 검색마다 맨 끝으로 재배치된다
        self._filter_trailing_widgets = [
            QLabel("날짜:"), self.filter_date_from, QLabel("~"),
            self.filter_date_to, filter_reset_btn,
        ]
        for w in self._filter_trailing_widgets:
            self.filter_ext_layout.addWidget(w)

        # 확장자가 많아 FlowLayout이 여러 줄로 줄바꿈될 수 있으므로 높이를 제한하지 않는다
        # (options2_container와 동일한 이유)
        layout.addWidget(filter_container)

        # ===== 탭: 결과 표시 =====
        self.tabs = QTabWidget()

        # 테이블 탭
        self.table = HoverableTableWidget()
        self.table.setColumnCount(5)
        self.table.setHorizontalHeaderLabels([
            "파일명",
            "경로",
            "크기",
            "수정 날짜",
            "검색 단어수"
        ])
        # 동적 크기 조정 설정
        self.table.horizontalHeader().setStretchLastSection(False)
        self.table.horizontalHeader().setSectionResizeMode(1, 0)  # 경로: 사용자가 드래그로 크기 조절 가능(Interactive)
        self.table.verticalHeader().setVisible(False)  # 행 번호 숨김
        self.table.setSelectionBehavior(0)  # 행 선택 모드
        # Qt 내장 정렬(setSortingEnabled)은 쓰지 않는다: 파일명/경로 컬럼은 아이콘·강조 표시를
        # cellWidget(QLabel)으로 그리는데 Qt 정렬은 QTableWidgetItem의 텍스트만 보고, 크기/검색
        # 단어수는 "1.2 MB"/"10"처럼 화면 표시용 문자열이라 문자열 정렬하면 숫자 순서와 어긋난다
        # (예: "10 KB"가 "2 MB"보다 앞에 옴). 대신 self.results(원본 값)를 정렬해 테이블을
        # 다시 그리는 방식으로 헤더 클릭 정렬을 구현한다.
        self.table.setSortingEnabled(False)
        self.table.setMouseTracking(True)  # ✅ 마우스 트래킹 활성화 (호버 감지용)
        self.table.cellDoubleClicked.connect(self.open_file)  # 더블클릭 시 파일 실행 (QLabel/Item 모두 처리)
        self.table.cellEntered.connect(self.on_table_cell_entered)  # 마우스 호버 시 미리보기 표시 (QLabel 포함)
        self.table.setContextMenuPolicy(Qt.CustomContextMenu)
        self.table.customContextMenuRequested.connect(self.show_table_context_menu)  # 우클릭 메뉴 (탐색기에서 위치 열기)
        self.table.horizontalHeader().setSortIndicatorShown(True)
        self.table.horizontalHeader().sectionClicked.connect(self.on_header_clicked)

        self.tabs.addTab(self.table, "🗂️ 결과 (테이블)")

        # 텍스트 탭
        self.text_output = QTextEdit()
        self.text_output.setReadOnly(True)
        font = QFont("Consolas")
        font.setPointSize(9)
        self.text_output.setFont(font)
        self.tabs.addTab(self.text_output, "📋 결과 (텍스트)")

        # 탭 스타일 설정
        self.tabs.setStyleSheet("""
            QTabWidget::pane {
                border: 1px solid #ddd;
            }
            QTabBar::tab {
                background-color: #f5f5f5;
                color: #333;
                padding: 6px 16px;
                margin-right: 2px;
                border: 1px solid #ddd;
                border-bottom: none;
                font-size: 10pt;
            }
            QTabBar::tab:selected {
                background-color: #ffffff;
                color: #0078d4;
                border: 1px solid #ddd;
                border-bottom: 3px solid #0078d4;
                font-weight: bold;
            }
            QTabBar::tab:hover:!selected {
                background-color: #efefef;
            }
        """)

        layout.addWidget(self.tabs, 1)

        # ===== 푸터 =====
        footer_container = QWidget()
        footer_layout = QHBoxLayout(footer_container)
        footer_layout.setContentsMargins(0, 0, 0, 0)
        self.result_count = QLabel("결과: 0개")
        footer_layout.addWidget(self.result_count)
        footer_layout.addStretch()
        footer_layout.addWidget(QLabel("fsearch v2.0 - GUI Edition"))
        footer_container.setMaximumHeight(20)
        layout.addWidget(footer_container)

    def update_excluded_stats_display(self):
        """자주 제외하는 폴더 통계 표시"""
        top_folders = self.search_history.get_top_excluded_folders(limit=3)
        stats = self.search_history.get_excluded_stats()

        if not top_folders:
            self.stats_label.setText("통계 없음")
            return

        # 상위 3개 폴더와 사용 횟수 표시
        stats_text = " | ".join(
            [f"{Path(folder).name} ({stats.get(folder, 0)}회)" for folder in top_folders]
        )
        self.stats_label.setText(stats_text)

    def add_exclude_folder(self):
        """제외할 폴더 추가"""
        folder = QFileDialog.getExistingDirectory(self, "제외할 폴더 선택")
        if folder:
            folder_name = Path(folder).name  # 폴더 이름만 추출
            full_path = folder  # 전체 경로 저장

            # 이미 추가되었는지 확인
            if full_path not in self.exclude_checkboxes:
                # 새로운 체크박스 생성
                cb = QCheckBox(folder_name)
                cb.setChecked(True)  # 새로 추가된 폴더는 제외됨
                cb.setMaximumHeight(25)
                cb.setToolTip(full_path)  # 전체 경로를 툴팁으로 표시

                # stretch 위젯 전에 체크박스 추가
                self.exclude_layout.insertWidget(
                    self.exclude_layout.count() - 1, cb
                )
                self.exclude_checkboxes[full_path] = cb

                # 제외 폴더 통계 기록
                self.search_history.add_excluded_folder(full_path)

                # 통계 표시 업데이트
                self.update_excluded_stats_display()

                self._set_status(f"✅ 제외 폴더 추가됨: {folder_name}")
            else:
                QMessageBox.information(self, "알림", "이미 추가된 폴더입니다.")

    def browse_path(self):
        """폴더 선택 대화창"""
        folder = QFileDialog.getExistingDirectory(self, "폴더 선택")
        if folder:
            self.path_input.setCurrentText(folder)
            self.search_history.add_path(folder)

    def closeEvent(self, event):
        """창을 닫을 때 검색이 진행 중이면 워커를 정리한 뒤 종료한다.
        기본 동작(오버라이드 없음)은 검색 중에도 창만 사라질 뿐 워커(및 그 안의
        ProcessPoolExecutor 자식 프로세스)는 백그라운드에서 계속 실행된다 - 실측:
        2,000개 파일을 처리하던 도중 닫아도 프로세스가 눈에 안 보이는 채로 7초간 더 실행됨."""
        if self.is_searching and self.search_worker and self.search_worker.isRunning():
            self.search_worker.stop_flag = True
            self._set_status("⏹️ 종료 중... 검색을 정리하고 있습니다.")
            # cancel_futures=True 덕분에 대부분 수 초 안에 끝난다(중지 버튼 실측 1.56초).
            # 그래도 안 끝나면 무한정 기다리지 않고 진행한다.
            if not self.search_worker.wait(5000):
                self.logger.warning("검색 워커가 5초 내에 종료되지 않아 강제로 앱을 닫습니다.")
        event.accept()

    def search(self):
        """검색 실행 또는 중지"""
        # 결과 초과 경고로 잠긴 상태면 검색어를 바꾸기 전까지 무시한다
        # (버튼을 setEnabled(False)로 실제 비활성화하면 Qt/Windows에서 disabled
        # 위젯이 hover 이벤트를 받지 않아 setCursor()로 지정한 커서가 표시되지
        # 않으므로, 버튼은 활성 상태로 두고 클릭 자체를 여기서 막는다)
        if getattr(self, '_search_locked', False):
            return
        # 검색 중이면 중지
        if self.is_searching:
            self.is_searching = False
            # 모든 타이머/동작 멈추기
            self.blink_timer.stop()
            self.set_status_urgent(False)
            self.progress_bar.setVisible(False)
            # 검색 워커 중단 신호
            if self.search_worker:
                self.search_worker.stop_flag = True
            # 상태 메시지 및 버튼 즉시 복원
            self._set_status("⏹️ 검색 중단됨")
            self.search_btn.setText("🔍 검색")
            self.search_btn.setStyleSheet("QPushButton { }")  # 기본 스타일 복원
            return

        keyword = self.keyword_input.currentText().strip()
        if not keyword:
            QMessageBox.warning(self, "입력 오류", "검색어를 입력하세요.")
            return

        path = self.path_input.currentText().strip()
        if not Path(path).exists():
            QMessageBox.warning(self, "경로 오류", "경로가 존재하지 않습니다.")
            return

        # 검색 이력에 저장 (최대 10개)
        self.search_history.add_keyword(keyword)

        # 검색 경로 저장 (최대 5개)
        self.search_history.add_path(path)

        # 현재 검색 단어 저장
        self.current_keyword = keyword
        self.current_regex = self.regex_cb.isChecked()

        # 검색 시작 시간 기록
        self.search_start_time = time.time()
        self.performance_btn.setText("⏱️ 완료시간 (진행중...)")
        self.performance_btn.setEnabled(False)

        # 스킵된 파일 갯수 초기화
        self.skipped_files_count_total = 0

        # 검색 시작 - 이전 결과 모두 초기화
        self.progress_bar.setVisible(True)
        self.progress_bar.setValue(0)
        self.table.setRowCount(0)
        self.text_output.clear()
        self.results = []

        # 이전 검색 결과 초기화
        self.excluded_files = []
        self.read_files = []
        self.no_match_files = []
        self.skipped_files = []

        # 관련 버튼 초기화
        self.excluded_btn.setText("❌ 제외된 파일 (0)")
        self.read_files_btn.setText("🔍 찾은 파일 수 (0)")
        self.no_match_files_btn.setText("❌ 검색어 미포함 파일 수 (0)")
        self.skipped_files_btn.setText("⏭️ 스킵된 대용량파일 (0)")

        # 체크된 제외 폴더만 수집
        ignore_dirs = {folder for folder, cb in self.exclude_checkboxes.items() if cb.isChecked()}

        self.search_worker = SearchWorker(
            keyword=keyword,
            path=path,
            ignore_dirs=list(ignore_dirs),
            name_only=self.name_only_cb.isChecked(),
            content_only=self.content_only_cb.isChecked(),
            use_regex=self.regex_cb.isChecked(),
            max_workers=self.workers_spin.value(),
            skip_large_files=self.skip_large_cb.isChecked()
        )

        self.search_worker.progress.connect(self.update_progress)
        self.search_worker.result_found.connect(self.add_result_row)  # 실시간 결과 추가
        self.search_worker.finished.connect(self.search_finished)
        self.search_worker.status.connect(self.update_status)
        self.search_worker.error.connect(self.search_error)
        self.search_worker.excluded_files_updated.connect(self.update_excluded_files)
        self.search_worker.file_processing.connect(self.update_current_file)  # 현재 처리 중인 파일
        self.search_worker.no_match_files_updated.connect(self.update_no_match_files)  # 검색어 미포함 파일
        self.search_worker.skipped_files_count.connect(self.update_skipped_files_count)  # 스킵된 파일 갯수
        self.search_worker.skipped_files_updated.connect(self.update_skipped_files)  # 스킵된 대용량파일 목록
        self.search_worker.extraction_urgent.connect(self.set_status_urgent)  # 캐시DB 생성 중 깜박임

        # 로깅
        self.logger.info(f"검색 시작 - 경로: {path}, 검색어: {keyword}")

        # 검색 시작 - 버튼 텍스트 변경 및 깜박임 시작
        self.is_searching = True
        self.search_btn.setText("⏹️ 중지")
        self.blink_state = False
        self.toggle_button_blink()  # 초기 상태 설정
        self.blink_timer.start(500)  # 500ms 간격으로 깜박임
        self.search_worker.start()

    def clear_search_history(self):
        """모든 검색어 삭제"""
        reply = QMessageBox.question(
            self,
            "검색어 초기화",
            "저장된 모든 검색어를 삭제하시겠습니까?",
            QMessageBox.Yes | QMessageBox.No,
            QMessageBox.No
        )
        if reply == QMessageBox.Yes:
            self.search_history.clear_keywords()
            self.keyword_input.clear()
            QMessageBox.information(self, "완료", "검색어가 모두 삭제되었습니다.")

    def clear_path_history(self):
        """모든 검색 경로 삭제"""
        reply = QMessageBox.question(
            self,
            "경로 초기화",
            "저장된 모든 검색 경로와 자주 제외하는 폴더 기록을 삭제하시겠습니까?",
            QMessageBox.Yes | QMessageBox.No,
            QMessageBox.No
        )
        if reply == QMessageBox.Yes:
            self.search_history.clear_paths()
            self.search_history.clear_excluded_stats()
            self.path_input.clear()
            self.update_excluded_stats_display()
            QMessageBox.information(self, "완료", "경로와 자주 제외하는 폴더 기록이 모두 삭제되었습니다.")

    def toggle_button_blink(self):
        """버튼 깜박임 토글"""
        self.blink_state = not self.blink_state
        if self.blink_state:
            # 깜박임 ON - 밝은 색상, 굵은 글씨, 큰 크기
            self.search_btn.setStyleSheet(
                "QPushButton { "
                "background-color: #FF6B6B; "
                "color: white; "
                "font-weight: bold; "
                "font-size: 12px; "
                "border: 2px solid #FF0000; "
                "border-radius: 4px; "
                "}"
            )
        else:
            # 깜박임 OFF - 기본 상태
            self.search_btn.setStyleSheet(
                "QPushButton { "
                "background-color: #FF6B6B; "
                "color: white; "
                "font-weight: bold; "
                "font-size: 12px; "
                "border: 2px solid #CC5555; "
                "border-radius: 4px; "
                "opacity: 0.5; "
                "}"
            )

    def set_status_urgent(self, urgent: bool):
        """캐시DB가 비어 있어 최초 대량 추출+색인 생성 중임을 상태표시바 깜박임으로 표시/해제"""
        if urgent:
            if not self.status_blink_timer.isActive():
                self.status_blink_state = False
                self.status_blink_timer.start(500)  # 500ms 간격 (버튼 깜박임과 동일 주기)
        else:
            self.status_blink_timer.stop()
            self.status_label.setStyleSheet("")

    def toggle_status_blink(self):
        """상태표시바 깜박임 토글 (붉은 굵은 글씨 <-> 기본)"""
        self.status_blink_state = not self.status_blink_state
        if self.status_blink_state:
            self.status_label.setStyleSheet("color: #FF0000; font-weight: bold;")
        else:
            self.status_label.setStyleSheet("")

    def on_keyword_changed(self, text):
        """결과 초과로 잠긴 검색을, 검색어를 다시 입력하면 풀어준다.
        init_ui()에서 검색 이력을 keyword_input에 addItems()로 채우는 시점에도
        이 시그널이 발생하는데, 그때는 아직 search_btn이 만들어지기 전이라 가드한다."""
        if not hasattr(self, 'search_btn'):
            return
        if getattr(self, '_search_locked', False):
            self._search_locked = False
            self.search_btn.setStyleSheet("QPushButton { }")
            self.search_btn.setCursor(Qt.ArrowCursor)
            self.status_label.setStyleSheet("")
            self._max_results_warning_active = False

    def _set_status(self, text):
        """상태표시바 텍스트를 갱신한다. 결과 초과 경고가 떠 있는 동안은 검색어를
        바꾸기 전까지(on_keyword_changed) 다른 이벤트가 이 문구를 덮어쓰지 못하게
        여기서 막는다. 경고 자체를 표시/해제하는 지점(search_finished)은 이 판단을
        내리는 쪽이므로 이 헬퍼를 거치지 않고 status_label을 직접 설정한다."""
        if getattr(self, '_max_results_warning_active', False):
            return
        self.status_label.setText(text)

    def update_progress(self, value):
        """진행바 업데이트"""
        self.progress_bar.setValue(value)

    def update_skipped_files_count(self, count):
        """스킵된 파일 갯수 업데이트"""
        self.skipped_files_count_total = count

    def add_result_row(self, result):
        """테이블에 결과 행 추가 (실시간)"""
        current_row_count = self.table.rowCount()
        self.table.insertRow(current_row_count)

        # 파일 크기 포맷
        size = result['size']
        if size < 1024:
            size_str = f"{size} B"
        elif size < 1024 * 1024:
            size_str = f"{size / 1024:.1f} KB"
        else:
            size_str = f"{size / (1024 * 1024):.1f} MB"

        # 각 컬럼에 데이터 입력
        # matched_lines를 미리 준비 (모든 셀에 저장)
        matched_lines = result.get('matched_lines', [])

        # 0: 파일명 (윈도우 탐색기 아이콘 + 검색어 강조 표시)
        keyword = getattr(self, 'current_keyword', '')
        use_regex = getattr(self, 'current_regex', False)
        filename_html = self._build_filename_display(result, keyword, use_regex)
        filename_label = QLabel()
        filename_label.setText(filename_html)
        filename_label.setStyleSheet("padding: 3px;")
        filename_label.setToolTip(result['full_path'])
        filename_label.setProperty("file_path", result['full_path'])  # ✅ 파일 경로 저장
        filename_label.setProperty("matched_lines", matched_lines)  # ✅ matched_lines 저장
        filename_label.setProperty("row", current_row_count)  # ✅ 행 번호 저장
        filename_label.setAttribute(Qt.WA_TransparentForMouseEvents)  # 마우스 이벤트를 테이블로 통과시켜 호버/더블클릭이 정상 동작하도록 함
        self.table.setCellWidget(current_row_count, 0, filename_label)
        # 위젯이 마우스 이벤트를 가로채지 않으므로, 툴팁/데이터는 동일 셀의 아이템에도 보관(다른 컬럼과 동일한 방식)
        filename_item = QTableWidgetItem("")
        filename_item.setToolTip(result['full_path'])
        filename_item.setData(Qt.UserRole, result['full_path'])
        filename_item.setData(Qt.UserRole + 1, matched_lines)
        self.table.setItem(current_row_count, 0, filename_item)

        # 1: 경로 (검색어 강조 표시)
        path_text = result['folder_path']
        if hasattr(self, 'current_keyword') and self.current_keyword:
            highlighted_path = self._highlight_keyword(path_text, self.current_keyword, self.current_regex)
            if '<span' in highlighted_path:
                path_label = QLabel()
                path_label.setText(highlighted_path)
                path_label.setStyleSheet("padding: 3px;")
                path_label.setToolTip(path_text)
                path_label.setProperty("file_path", result['full_path'])  # ✅ 파일 경로 저장
                path_label.setProperty("matched_lines", matched_lines)  # ✅ matched_lines 저장
                path_label.setProperty("row", current_row_count)  # ✅ 행 번호 저장
                self.table.setCellWidget(current_row_count, 1, path_label)
            else:
                path_item = QTableWidgetItem(path_text)
                path_item.setToolTip(path_text)
                path_item.setTextAlignment(Qt.AlignLeft | Qt.AlignVCenter)
                path_item.setData(Qt.UserRole, result['full_path'])  # ✅ 파일 경로 저장
                path_item.setData(Qt.UserRole + 1, matched_lines)  # ✅ matched_lines 저장
                self.table.setItem(current_row_count, 1, path_item)
        else:
            path_item = QTableWidgetItem(path_text)
            path_item.setToolTip(path_text)
            path_item.setTextAlignment(Qt.AlignLeft | Qt.AlignVCenter)
            path_item.setData(Qt.UserRole, result['full_path'])  # ✅ 파일 경로 저장
            path_item.setData(Qt.UserRole + 1, matched_lines)  # ✅ matched_lines 저장
            self.table.setItem(current_row_count, 1, path_item)

        # 2: 크기 (오른쪽 정렬 - 숫자)
        size_item = QTableWidgetItem(size_str)
        size_item.setToolTip(str(result['size']) + " bytes")
        size_item.setTextAlignment(Qt.AlignRight | Qt.AlignVCenter)
        size_item.setData(Qt.UserRole, result['full_path'])  # ✅ 파일 경로 저장
        # matched_lines를 UserRole+1로 저장 (미리보기용)
        matched_lines = result.get('matched_lines', [])
        size_item.setData(Qt.UserRole + 1, matched_lines)
        self.table.setItem(current_row_count, 2, size_item)

        # 3: 수정 날짜 (왼쪽 정렬)
        modified_item = QTableWidgetItem(result['modified'])
        modified_item.setTextAlignment(Qt.AlignLeft | Qt.AlignVCenter)
        modified_item.setData(Qt.UserRole, result['full_path'])  # ✅ 파일 경로 저장
        modified_item.setData(Qt.UserRole + 1, matched_lines)  # ✅ matched_lines 저장
        self.table.setItem(current_row_count, 3, modified_item)

        # 4: 검색 단어수 (오른쪽 정렬 - 숫자)
        match_count = result.get('match_count', 0)
        match_count_item = QTableWidgetItem(str(match_count))
        match_count_item.setTextAlignment(Qt.AlignRight | Qt.AlignVCenter)
        match_count_item.setData(Qt.UserRole, result['full_path'])  # ✅ 파일 경로 저장
        match_count_item.setData(Qt.UserRole + 1, matched_lines)  # ✅ matched_lines 저장
        self.table.setItem(current_row_count, 4, match_count_item)

    def update_status(self, status):
        """상태 메시지 업데이트"""
        self._set_status(status)

    def update_current_file(self, file_info):
        """현재 처리 중인 파일 표시"""
        self._set_status(f"검색 중: {file_info}")

    def update_no_match_files(self, no_match_files):
        """검색어 미포함 파일 누적 업데이트"""
        # 검색 시작 시 초기화되고 검색당 한 번만 호출되며 입력 자체도 중복이 없어 중복 체크 불필요
        self.no_match_files.extend(no_match_files)

        # 버튼 텍스트 업데이트
        self.no_match_files_btn.setText(f"❌ 검색어 미포함 파일 수 ({len(self.no_match_files)})")

    def update_skipped_files(self, skipped_files):
        """스킵된 대용량파일 목록 누적 업데이트"""
        # 검색 시작 시 초기화되고 검색당 한 번만 호출되며 입력 자체도 중복이 없어 중복 체크 불필요
        self.skipped_files.extend(skipped_files)

    def search_finished(self, results):
        """검색 완료 또는 중단"""
        # 이미 사용자가 중지 버튼을 눌렀으면 UI는 이미 복원됨
        if not self.is_searching:
            return

        # 검색 완료 (중단되지 않음)
        # 완료 시간 계산
        if self.search_start_time:
            self.search_elapsed_time = time.time() - self.search_start_time
            self.performance_btn.setText(f"⏱️ 완료시간 ({self.search_elapsed_time:.2f}초)")
            self.performance_btn.setEnabled(True)

        # 검색 상태 복원
        self.is_searching = False
        self.blink_timer.stop()
        self.set_status_urgent(False)
        self.search_btn.setText("🔍 검색")
        self.search_btn.setStyleSheet("QPushButton { }")  # 기본 스타일 복원

        self.results = results
        self.progress_bar.setVisible(False)

        # 검색 단어를 포함하는 셀을 굵게 표시 + 빨간색
        if hasattr(self, 'current_keyword'):
            keyword = self.current_keyword
            keyword_query = None if self.current_regex else _parse_keyword_query(keyword)

            for row in range(self.table.rowCount()):
                # 파일명 (컬럼 0)과 경로 (컬럼 1) 체크
                for col in [0, 1]:
                    item = self.table.item(row, col)
                    if item:
                        text = item.text()
                        # 검색 단어 찾기
                        found = False
                        if self.current_regex:
                            try:
                                if re.search(keyword, text, re.IGNORECASE):
                                    found = True
                            except:
                                pass
                        else:
                            if _query_matches(keyword_query, text.lower()):
                                found = True

                        # 찾으면 bold 글꼴 + 빨간색 적용
                        if found:
                            bold_font = QFont(item.font())
                            bold_font.setBold(True)
                            item.setFont(bold_font)
                            item.setForeground(QColor('red'))

        # 검색 단어수로 정렬
        self.sort_results_by_match_count()

        # 필터 UI를 새 결과의 확장자 집합으로 재구성(전부 체크된 상태로 시작)하고,
        # 테이블/텍스트 탭을 그린다(날짜 필터는 이전 검색값을 유지한 채 적용됨)
        self._rebuild_extension_filter()
        self._render_results()

        # 파일별 검색 결과 개수 계산 (필터와 무관하게 이번 검색의 전체 결과 기준)
        file_counts = defaultdict(int)
        for result in results:
            file_counts[result['full_path']] += 1

        # 읽은 파일 누적 (검색 시작 시 초기화되고 검색당 한 번만 호출되며 키는 이미 유일하므로 중복 체크 불필요)
        self.read_files.extend(file_counts.keys())

        # 읽은 파일 버튼 업데이트
        self.read_files_btn.setText(f"🔍 찾은 파일 수 ({len(self.read_files)})")

        # 스킵된 대용량파일 버튼 업데이트
        self.skipped_files_btn.setText(f"⏭️ 스킵된 대용량파일 ({self.skipped_files_count_total})")
        self.skipped_files_btn.setEnabled(self.skipped_files_count_total > 0)

        # 상태 메시지 업데이트
        # MAX_RESULTS 초과로 중단된 경우, worker가 emit한 경고 상태 메시지가
        # 이 "완료" 메시지로 즉시 덮어써져 사용자가 볼 새도 없이 사라지므로
        # (finished 신호 직전에 emit되지만 GUI 스레드에서 연속 처리됨) 유지한다.
        if self.search_worker and getattr(self.search_worker, '_max_results_hit', False):
            self.status_label.setStyleSheet("color: #00008B; font-weight: bold;")
            self.status_label.setText(
                f"⚠️ 결과가 {self.search_worker.MAX_RESULTS}건을 초과해 검색을 중단했습니다 "
                f"({len(results)}개 표시). 검색어를 더 구체적으로 입력해주세요."
            )
            # 검색어를 바꾸기 전까지는 같은 상황이 반복될 뿐이므로 검색을 막는다
            # (on_keyword_changed가 검색어 입력칸 텍스트 변경 시 다시 풀어준다)
            self._search_locked = True
            self.search_btn.setStyleSheet("QPushButton { color: gray; font-weight: bold; }")
            self.search_btn.setCursor(Qt.ForbiddenCursor)
            # 결과 테이블에 마우스를 올리면 on_table_cell_entered가 미리보기 문구로
            # 상태바를 덮어써 경고가 곧바로 사라지므로, 검색어를 바꾸기 전까지 막는다
            self._max_results_warning_active = True
        else:
            self.status_label.setStyleSheet("")
            self.status_label.setText(f"✅ 검색 완료: {len(results)}개 결과")
            self._max_results_warning_active = False

        # 로깅 - 검색 결과
        self.logger.info(f"검색 완료 - 총 {len(results)}개 결과 (파일: {len(file_counts)}개, 스킵: {self.skipped_files_count_total}개)")
        if self.search_worker and hasattr(self.search_worker, 'timing'):
            t = self.search_worker.timing
            self.logger.info(
                f"구간별 시간 - 수집:{t.get('collect', 0):.3f}s "
                f"캐시로드:{t.get('cache_load', 0):.3f}s "
                f"캐시정리:{t.get('cache_cleanup', 0):.3f}s "
                f"검색:{t.get('search', 0):.3f}s "
                f"분류:{t.get('classify', 0):.3f}s "
                f"합계:{t.get('total', 0):.3f}s"
            )
            self.logger.info(
                f"검색구간 진행률 체크포인트 - "
                f"25%:{t.get('checkpoint_25pct', 0):.3f}s "
                f"50%:{t.get('checkpoint_50pct', 0):.3f}s "
                f"75%:{t.get('checkpoint_75pct', 0):.3f}s "
                f"100%:{t.get('checkpoint_100pct', 0):.3f}s"
            )

    def _rebuild_extension_filter(self):
        """새 검색 결과의 확장자 집합으로 필터 체크박스를 다시 만든다(전부 체크된 상태로 시작).
        날짜/초기화 위젯은 없애지 않고 체크박스 뒤로 다시 옮겨 붙여, 항상 마지막
        체크박스 바로 옆에 이어지도록 한다."""
        for cb in self.filter_ext_checkboxes.values():
            cb.setParent(None)
        self.filter_ext_checkboxes = {}

        for w in self._filter_trailing_widgets:
            self.filter_ext_layout.removeWidget(w)

        extensions = sorted({r.get('extension', '') or '(확장자 없음)' for r in self.results})
        for ext in extensions:
            cb = QCheckBox(ext)
            cb.setChecked(True)
            cb.setMaximumHeight(25)
            cb.stateChanged.connect(self._apply_filters)
            self.filter_ext_layout.addWidget(cb)
            self.filter_ext_checkboxes[ext] = cb

        for w in self._filter_trailing_widgets:
            self.filter_ext_layout.addWidget(w)

    def reset_filters(self):
        """확장자 필터를 전부 체크, 날짜 필터를 비운 뒤 다시 그린다"""
        for cb in self.filter_ext_checkboxes.values():
            cb.blockSignals(True)
            cb.setChecked(True)
            cb.blockSignals(False)
        self.filter_date_from.blockSignals(True)
        self.filter_date_from.clear()
        self.filter_date_from.blockSignals(False)
        self.filter_date_to.blockSignals(True)
        self.filter_date_to.clear()
        self.filter_date_to.blockSignals(False)
        self._filter_debounce_timer.stop()
        self._apply_filters()

    @staticmethod
    def _parse_filter_date(text):
        """YYYY-MM-DD 형식이 아니거나 비어 있으면 None(필터 미적용)"""
        text = text.strip()
        if not text:
            return None
        try:
            return datetime.strptime(text, '%Y-%m-%d')
        except ValueError:
            return None

    @staticmethod
    def _parse_result_date(text):
        """result['modified']('%Y-%m-%d %H:%M:%S' 또는 stat 실패 시 'Unknown')를 파싱"""
        try:
            return datetime.strptime(text, '%Y-%m-%d %H:%M:%S')
        except (ValueError, TypeError):
            return None

    def _result_passes_filters(self, result) -> bool:
        """확장자/날짜 필터를 모두 통과하면 True"""
        ext = result.get('extension', '') or '(확장자 없음)'
        cb = self.filter_ext_checkboxes.get(ext)
        if cb and not cb.isChecked():
            return False

        date_from = self._parse_filter_date(self.filter_date_from.text())
        date_to = self._parse_filter_date(self.filter_date_to.text())
        if date_from or date_to:
            modified = self._parse_result_date(result.get('modified', ''))
            if modified is None:
                return False  # 날짜를 알 수 없는 결과는 날짜 필터가 걸려 있으면 제외
            if date_from and modified.date() < date_from.date():
                return False
            if date_to and modified.date() > date_to.date():
                return False

        return True

    def _schedule_apply_filters(self):
        """날짜 입력칸에서 타이핑이 잠깐(300ms) 멈추면 그때 한 번만 필터를 적용한다
        (매 키 입력마다 적용하면 결과가 많을 때 입력이 밀리는 것처럼 느껴진다)"""
        self._filter_debounce_timer.start(300)

    def _render_results(self):
        """self.results 전체로 테이블 행을 새로 만든다(위젯을 새로 생성하는 무거운 동작).
        행 순서 자체가 바뀔 때만(검색 완료 직후, 헤더 클릭 정렬 후) 호출한다.
        다 그린 뒤 현재 필터 상태를 이어서 반영한다(_apply_filters).
        결과가 수천 건일 때 필터를 바꿀 때마다 이 함수를 다시 부르면 매번 위젯을 전부
        새로 만들어 화면이 멈춘 것처럼 느려지므로, 필터 토글은 이 함수를 거치지 않고
        _apply_filters가 행을 숨기는 것만으로 처리한다."""
        self.table.setRowCount(0)
        for result in self.results:
            self.add_result_row(result)
        self.table.resizeColumnsToContents()
        self._apply_filters()

    def _apply_filters(self):
        """확장자/날짜 필터를 반영한다. 테이블은 이미 그려진 행을 숨기고 보이기만 해서
        (setRowHidden - 위젯을 새로 만들지 않아 가볍다) 반영하고, 텍스트 탭/결과 수는
        위젯 생성이 없는 가벼운 작업이라 필터된 목록으로 매번 다시 그려도 무리 없다."""
        filtered = []
        for row, result in enumerate(self.results):
            passes = self._result_passes_filters(result)
            self.table.setRowHidden(row, not passes)
            if passes:
                filtered.append(result)
        self._render_text_tab(filtered)

    def _render_text_tab(self, results):
        """텍스트 탭과 결과 수 표시를 그린다(위젯 생성이 없어 필터를 토글할 때마다 불러도 가볍다)"""
        # 텍스트 탭 업데이트 (검색어 굵게 표기)
        text_output = "<html><body><pre>검색 결과:\n" + "="*100 + "\n\n"
        shown_files = set()

        # 검색어 준비
        keyword = getattr(self, 'current_keyword', '')
        use_regex = getattr(self, 'current_regex', False)

        for result in results:
            # 각 파일당 한 번씩만 정보 표시
            file_path = result['full_path']
            if file_path not in shown_files:
                shown_files.add(file_path)

                size = result['size']
                if size < 1024:
                    size_str = f"{size} B"
                elif size < 1024 * 1024:
                    size_str = f"{size / 1024:.1f} KB"
                else:
                    size_str = f"{size / (1024 * 1024):.1f} MB"

                # 파일명에서 검색어를 굵게 표기
                filename = result['filename']
                highlighted_filename = self._highlight_keyword(filename, keyword, use_regex)

                # 경로에서 검색어를 굵게 표기
                highlighted_path = self._highlight_keyword(file_path, keyword, use_regex)

                text_output += f"[{highlighted_filename}]\n"
                text_output += f"  경로: {highlighted_path}\n"
                text_output += f"  크기: {size_str}\n"
                text_output += f"  수정일: {result['modified']}\n"
                text_output += f"  검색 단어수: {result['match_count']}\n"

                # 매칭된 라인 출력 (최대 3줄)
                matched_lines = result.get('matched_lines', [])
                if matched_lines:
                    text_output += f"  검색어 포함 문장:\n"
                    for idx, line in enumerate(matched_lines, 1):
                        # 검색어 강조
                        highlighted_line = self._highlight_keyword(line, keyword, use_regex)
                        # 긴 라인은 줄임
                        if len(highlighted_line) > 120:
                            highlighted_line = highlighted_line[:120] + "..."
                        text_output += f"    {idx}. {highlighted_line}\n"

                text_output += "\n"

        text_output += "</pre></body></html>"
        self.text_output.setHtml(text_output)

        # 파일별 표시 결과 개수 계산(필터가 적용된 현재 표시 기준)
        file_counts = defaultdict(int)
        for result in results:
            file_counts[result['full_path']] += 1

        self.result_count.setText(f"결과: {len(results)}개 (파일: {len(file_counts)}개)")

    def _highlight_keyword(self, text: str, keyword: str, use_regex: bool) -> str:
        """텍스트에서 검색어를 HTML 굵게 + 빨간색으로 표기"""
        if not keyword:
            return text

        if use_regex:
            try:
                # 정규식 사용
                pattern = re.compile(f'({keyword})', re.IGNORECASE)
                return pattern.sub(r'<span style="color: red; font-weight: bold;">\1</span>', text)
            except:
                return text
        else:
            # 일반 문자열 검색 (대소문자 무시, AND/OR 다중 키워드는 항별로 각각 강조)
            # 모든 일치 부분을 굵게 + 빨간색으로 표기
            terms = sorted({t for group in _parse_keyword_query(keyword) for t in group}, key=len, reverse=True)
            if not terms:
                return text
            pattern = re.compile('(' + '|'.join(re.escape(t) for t in terms) + ')', re.IGNORECASE)
            return pattern.sub(r'<span style="color: red; font-weight: bold;">\1</span>', text)

    def _prewarm_icon_cache(self):
        """지원 확장자의 아이콘을 미리 조회해 캐싱 (검색 중 첫 조회 지연 방지)"""
        for ext in SearchWorker.EXTRACTABLE_EXTENSIONS | {'.zip'}:
            self._get_icon_html(ext)

    def _get_icon_html(self, extension: str) -> str:
        """확장자에 대해 윈도우 탐색기와 동일한 아이콘을 <img> HTML로 반환 (캐시 사용)"""
        if extension not in self._icon_html_cache:
            icon = self._icon_provider.icon(QFileInfo(f"dummy{extension}"))
            pixmap = icon.pixmap(16, 16)
            buf = QByteArray()
            qbuf = QBuffer(buf)
            qbuf.open(QIODevice.WriteOnly)
            pixmap.save(qbuf, "PNG")
            b64 = base64.b64encode(buf.data()).decode('ascii')
            self._icon_html_cache[extension] = (
                f'<img src="data:image/png;base64,{b64}" width="16" height="16" '
                'style="vertical-align: middle;">'
            )
        return self._icon_html_cache[extension]

    def _build_filename_display(self, result, keyword: str, use_regex: bool) -> str:
        """파일명 셀에 표시할 HTML(윈도우 탐색기 아이콘 + 검색어 강조 텍스트) 생성"""
        full_path = result['full_path']
        if 'inner_name' in result:
            # zip 내부 파일: 압축파일 아이콘 → 내부파일 아이콘
            text = f"{Path(full_path).name} → {result['inner_name']}"
            icon_html = f"{self._get_icon_html('.zip')} {self._get_icon_html(result['extension'])}"
        else:
            text = Path(full_path).name
            icon_html = self._get_icon_html(result['extension'])

        display_text = self._highlight_keyword(text, keyword, use_regex) if keyword else text
        return f"{icon_html} {display_text}"

    def search_error(self, error):
        """검색 오류"""
        self.progress_bar.setVisible(False)
        # 검색 상태 복원
        self.is_searching = False
        self.blink_timer.stop()
        self.set_status_urgent(False)
        self.search_btn.setText("🔍 검색")
        self.search_btn.setStyleSheet("QPushButton { }")  # 기본 스타일 복원
        QMessageBox.critical(self, "오류", error)
        self._set_status("오류 발생")
        # 로깅
        self.logger.error(f"검색 오류: {error}")

    def on_table_cell_entered(self, row, column):
        """테이블 셀 호버 시 미리보기 표시 및 행 하이라이트 (QLabel 위젯과 QTableWidgetItem 모두 처리)"""
        if row < 0 or row >= self.table.rowCount():
            return

        # 이전 호버 행의 배경색 복원
        if self.table.last_hovered_row >= 0 and self.table.last_hovered_row != row:
            for col in range(self.table.columnCount()):
                # QTableWidgetItem 처리
                item = self.table.item(self.table.last_hovered_row, col)
                if item:
                    item.setBackground(QColor(255, 255, 255))  # 흰색으로 복원
                # QLabel 위젯 처리
                widget = self.table.cellWidget(self.table.last_hovered_row, col)
                if widget and isinstance(widget, QLabel):
                    widget.setStyleSheet("background-color: white;")

        # 현재 행의 배경색을 진한 그레이로 설정
        gray_color = QColor(200, 200, 200)  # 진한 그레이 색상
        for col in range(self.table.columnCount()):
            # QTableWidgetItem 처리
            item = self.table.item(row, col)
            if item:
                item.setBackground(gray_color)
            # QLabel 위젯 처리
            widget = self.table.cellWidget(row, col)
            if widget and isinstance(widget, QLabel):
                widget.setStyleSheet("background-color: rgb(200, 200, 200);")

        # 현재 호버 행 저장
        self.table.last_hovered_row = row

        # matched_lines 가져오기 (여러 소스에서 시도)
        matched_lines = None
        found_source = None
        match_count = 0

        # 1. self.results에서 먼저 확인
        if row < len(self.results):
            matched_lines = self.results[row].get('matched_lines', [])
            match_count = self.results[row].get('match_count', 0)
            if matched_lines:
                found_source = "self.results"

        # 2. 테이블 아이템(QTableWidgetItem)에서 확인
        if not matched_lines:
            item = self.table.item(row, column)
            if item:
                matched_lines = item.data(Qt.UserRole + 1)

        # 3. QLabel 위젯의 property에서 확인
        if not matched_lines:
            widget = self.table.cellWidget(row, column)
            if widget and isinstance(widget, QLabel):
                matched_lines = widget.property("matched_lines")

        # 4. 모든 컬럼의 QLabel에서 확인
        if not matched_lines:
            for col in range(self.table.columnCount()):
                widget = self.table.cellWidget(row, col)
                if widget and isinstance(widget, QLabel):
                    matched_lines = widget.property("matched_lines")
                    if matched_lines:
                        break

        # 파일명 전용 검색 확인: match_count가 1이고 matched_lines가 비어있음
        if match_count == 1 and (not matched_lines or len(matched_lines) == 0):
            # 파일명에만 검색어가 있는 경우
            preview_text = "📝 본문에 검색어 없음\n" + "=" * 70 + "\n검색어는 파일명에만 포함되어 있습니다."

            # 해당 셀에 ToolTip 설정
            item = self.table.item(row, column)
            if item:
                item.setToolTip(preview_text)
            else:
                # QLabel 위젯인 경우
                widget = self.table.cellWidget(row, column)
                if widget and isinstance(widget, QLabel):
                    widget.setToolTip(preview_text)

            # 상태 바에도 메시지 표시
            self._set_status("📝 본문에 검색어 없음 (파일명만 일치)")
            return

        # matched_lines가 있으면 미리보기 표시
        if matched_lines and isinstance(matched_lines, list) and len(matched_lines) > 0:
            # ToolTip 텍스트 생성
            preview_text = "📝 검색어 포함 문장:\n" + "=" * 70 + "\n"

            for idx, line in enumerate(matched_lines, 1):
                display_line = line if len(line) <= 100 else line[:100] + "..."
                preview_text += f"{idx}. {display_line}\n"

            # 해당 셀에 ToolTip 설정
            item = self.table.item(row, column)
            if item:
                item.setToolTip(preview_text)
            else:
                # QLabel 위젯인 경우
                widget = self.table.cellWidget(row, column)
                if widget and isinstance(widget, QLabel):
                    widget.setToolTip(preview_text)

            # 상태 바에도 미리보기 표시 (첫 번째 문장)
            if matched_lines and len(matched_lines) > 0:
                status_msg = matched_lines[0]
                if len(status_msg) > 100:
                    status_msg = status_msg[:100] + "..."
                self._set_status(f"📝 {status_msg}")

    def open_file(self, row, column):
        """파일 또는 폴더 열기 (더블클릭 시 호출 - 검색 중/완료 모두 처리)"""
        try:
            # 행 번호 범위 검증
            if row < 0 or row >= self.table.rowCount():
                self._set_status("❌ 유효하지 않은 행입니다")
                return

            file_path = None

            # 1️⃣ self.results에서 파일 경로 추출 시도 (정렬되지 않은 상태에서만 가능)
            if row < len(self.results):
                result = self.results[row]
                file_path = result.get('full_path') or result.get('path')

            # 2️⃣ self.results에 없으면 테이블의 셀에서 추출
            # - QTableWidgetItem의 UserRole 데이터
            # - QLabel의 property
            if not file_path:
                # 크기 셀(2)에서 추출 시도
                item = self.table.item(row, 2)
                if item:
                    file_path = item.data(Qt.UserRole)

                # 여전히 없으면 QLabel에서 추출
                if not file_path:
                    widget = self.table.cellWidget(row, 0)
                    if widget and isinstance(widget, QLabel):
                        file_path = widget.property("file_path")

                    if not file_path:
                        widget = self.table.cellWidget(row, 1)
                        if widget and isinstance(widget, QLabel):
                            file_path = widget.property("file_path")

            # 파일 경로 검증
            if not file_path or not isinstance(file_path, str):
                self._set_status("❌ 파일 경로 정보가 없습니다")
                return

            # 문자열 정리
            file_path = str(file_path).strip()
            if not file_path:
                self._set_status("❌ 파일 경로가 비어있습니다")
                return

            # 경로 존재 여부 확인
            path_obj = Path(file_path)
            if not path_obj.exists():
                self._set_status(f"❌ 파일을 찾을 수 없습니다: {file_path}")
                return

            # 파일 실행
            if sys.platform == 'win32':
                os.startfile(file_path)
            elif sys.platform == 'darwin':
                subprocess.Popen(['open', file_path])
            else:
                subprocess.Popen(['xdg-open', file_path])

            filename = path_obj.name
            self._set_status(f"✅ 파일 열음: {filename}")

            # 로깅
            if hasattr(self, 'logger'):
                self.logger.info(f"파일 실행 - {file_path}")

        except Exception as e:
            error_msg = f"파일 실행 오류: {str(e)}"
            self._set_status(f"❌ {error_msg}")
            if hasattr(self, 'logger'):
                self.logger.error(error_msg)

    def _resolve_row_file_path(self, row):
        """행 번호로부터 파일 경로를 찾는다 (open_file과 동일한 다단계 폴백:
        self.results → 크기 셀 UserRole → 파일명/경로 셀 QLabel property)"""
        if row < 0 or row >= self.table.rowCount():
            return None

        file_path = None
        if row < len(self.results):
            result = self.results[row]
            file_path = result.get('full_path') or result.get('path')

        if not file_path:
            item = self.table.item(row, 2)
            if item:
                file_path = item.data(Qt.UserRole)

            if not file_path:
                widget = self.table.cellWidget(row, 0)
                if widget and isinstance(widget, QLabel):
                    file_path = widget.property("file_path")

                if not file_path:
                    widget = self.table.cellWidget(row, 1)
                    if widget and isinstance(widget, QLabel):
                        file_path = widget.property("file_path")

        if not file_path or not isinstance(file_path, str):
            return None

        return str(file_path).strip() or None

    def show_table_context_menu(self, pos):
        """결과 테이블 우클릭 메뉴 - 탐색기에서 파일 위치 열기"""
        row = self.table.rowAt(pos.y())
        file_path = self._resolve_row_file_path(row)
        if not file_path or not Path(file_path).exists():
            return

        menu = QMenu(self)
        action = menu.addAction("📂 파일 위치 열기")
        action.triggered.connect(lambda: self._show_in_explorer(file_path))
        menu.exec_(self.table.viewport().mapToGlobal(pos))

    def _show_in_explorer(self, file_path):
        """탐색기를 열어 해당 파일을 선택된 상태로 보여준다 (zip 내부 파일은 zip 자체를 선택)"""
        try:
            if sys.platform == 'win32':
                subprocess.Popen(f'explorer /select,"{file_path}"')
            else:
                subprocess.Popen(['xdg-open', str(Path(file_path).parent)])
            self._set_status(f"📂 탐색기에서 위치 열음: {Path(file_path).name}")
        except Exception as e:
            self._set_status(f"❌ 탐색기 열기 실패: {e}")
            if hasattr(self, 'logger'):
                self.logger.error(f"탐색기 열기 실패({file_path}): {e}")

    def calculate_relevance_score(self, result):
        """파일의 업무 연관도 점수 계산"""
        keyword = self.current_keyword if hasattr(self, 'current_keyword') else ""
        filename = result.get('filename', '').lower()
        extension = result.get('extension', '').lower()
        match_count = result.get('match_count', 0)

        score = 0

        # 1. 파일명에 검색어 포함: 최우선
        if keyword:
            use_regex = getattr(self, 'current_regex', False)
            if use_regex:
                if keyword.lower() in filename:
                    score += 1000
            elif _query_matches(_parse_keyword_query(keyword), filename):
                score += 1000

        # 2. 문서 타입 우선순위
        doc_extensions = {'.doc': 50, '.docx': 50, '.pdf': 40, '.txt': 30,
                         '.md': 30, '.xlsx': 40, '.xls': 40, '.hwp': 50}
        score += doc_extensions.get(extension, 0)

        # 3. 소스코드 제외 (낮은 우선순위)
        code_extensions = {'.py', '.js', '.cpp', '.java', '.cs', '.go', '.rs'}
        if extension in code_extensions:
            score -= 10

        # 4. match_count 기반 (기본 정렬 기준)
        score += match_count * 10

        return score

    # 헤더 클릭 정렬의 컬럼별 정렬 키. 파일명은 zip 내부 표시("압축파일 → 내부파일")까지 반영해
    # 엑셀 내보내기(export_to_excel)와 동일한 이름 규칙을 쓴다.
    def _sort_key_filename(self, result):
        name = Path(result['full_path']).name
        if 'inner_name' in result:
            name = f"{name} → {result['inner_name']}"
        return name.lower()

    _HEADER_SORT_KEYS = {
        0: _sort_key_filename,
        1: lambda self, result: result['folder_path'].lower(),
        2: lambda self, result: result['size'],
        3: lambda self, result: result['modified'],  # "YYYY-MM-DD HH:MM:SS" 형식이라 문자열 정렬 = 시간순
        4: lambda self, result: result.get('match_count', 0),
    }

    def on_header_clicked(self, column):
        """테이블 헤더 클릭 시 해당 컬럼 기준으로 정렬하고 테이블을 다시 그린다.
        같은 컬럼을 다시 클릭하면 오름차순/내림차순을 토글한다."""
        if column not in self._HEADER_SORT_KEYS:
            return
        if self.sort_column == column:
            self.sort_ascending = not self.sort_ascending
        else:
            self.sort_column = column
            self.sort_ascending = True

        key_fn = self._HEADER_SORT_KEYS[column]
        self.results.sort(key=lambda r: key_fn(self, r), reverse=not self.sort_ascending)

        # 정렬로 행 순서 자체가 바뀌므로 테이블을 다시 그린다(그 안에서 필터도 다시 반영됨)
        self._render_results()

        order = Qt.AscendingOrder if self.sort_ascending else Qt.DescendingOrder
        self.table.horizontalHeader().setSortIndicator(column, order)

    def sort_results_by_match_count(self):
        """업무 연관도 기준으로 self.results를 정렬한다(테이블 반영은 호출부의 _render_results가 담당)"""
        self.sort_column = None  # 새 검색 결과이므로 이전 헤더 클릭 정렬 기준을 초기화
        # self.results를 연관도 점수 기준 내림차순으로 정렬
        self.results.sort(key=lambda x: self.calculate_relevance_score(x), reverse=True)

    def export_to_excel(self):
        """현재 결과 테이블에 표시된 내용을 엑셀 파일로 내보내기"""
        if not Workbook:
            QMessageBox.warning(self, "엑셀 내보내기", "openpyxl 라이브러리가 설치되어 있지 않습니다.")
            return

        if not self.results:
            QMessageBox.information(self, "엑셀 내보내기", "내보낼 검색 결과가 없습니다.")
            return

        file_path, _ = QFileDialog.getSaveFileName(
            self, "엑셀로 저장", "fsearch_결과.xlsx", "Excel 파일 (*.xlsx)"
        )
        if not file_path:
            return

        try:
            wb = Workbook()
            ws = wb.active
            ws.title = "검색 결과"
            ws.append(["파일명", "경로", "크기", "수정 날짜", "검색 단어수"])

            for result in self.results:
                size = result['size']
                if size < 1024:
                    size_str = f"{size} B"
                elif size < 1024 * 1024:
                    size_str = f"{size / 1024:.1f} KB"
                else:
                    size_str = f"{size / (1024 * 1024):.1f} MB"

                filename = Path(result['full_path']).name
                if 'inner_name' in result:
                    filename = f"{filename} → {result['inner_name']}"

                ws.append([
                    filename,
                    result['folder_path'],
                    size_str,
                    result['modified'],
                    result.get('match_count', 0)
                ])

            wb.save(file_path)
            QMessageBox.information(self, "엑셀 내보내기", f"저장 완료:\n{file_path}")
        except Exception as e:
            QMessageBox.critical(self, "엑셀 내보내기 오류", f"저장 중 오류가 발생했습니다:\n{e}")

    def update_excluded_files(self, excluded_files):
        """제외된 파일 목록 업데이트"""
        self.excluded_files = excluded_files
        # 버튼 텍스트에 제외된 파일 개수 표시
        count = len(excluded_files)
        self.excluded_btn.setText(f"📁 제외된 파일 ({count})")

    def show_excluded_files(self):
        """제외된 파일 목록 보기"""
        if not self.excluded_files:
            QMessageBox.information(self, "제외된 파일", "제외된 파일이 없습니다.\n\n먼저 검색을 실행해주세요.")
            return

        # 제외된 파일 목록을 보여주는 윈도우
        from PyQt5.QtWidgets import QDialog
        dialog = QDialog(self)
        dialog.setWindowTitle("📁 제외된 파일 목록")
        dialog.setGeometry(200, 200, 800, 600)

        layout = QVBoxLayout(dialog)

        # 통계 정보
        info_label = QLabel(f"총 {len(self.excluded_files)}개의 파일이 제외되었습니다.")
        layout.addWidget(info_label)

        # 제외된 파일 목록
        text_edit = QTextEdit()
        text_edit.setReadOnly(True)
        text_edit.setFont(QFont("Consolas", 9))

        excluded_text = "\n".join(self.excluded_files)
        text_edit.setText(excluded_text)
        layout.addWidget(text_edit)

        # 닫기 버튼
        close_btn = QPushButton("닫기")
        close_btn.clicked.connect(dialog.close)
        layout.addWidget(close_btn)

        dialog.exec_()

    def show_read_files(self):
        """읽은 파일 목록 보기"""
        if not self.read_files:
            QMessageBox.information(self, "찾은 파일 수", "찾은 파일이 없습니다.\n\n먼저 검색을 실행해주세요.")
            return

        # 읽은 파일 목록을 보여주는 윈도우
        from PyQt5.QtWidgets import QDialog
        dialog = QDialog(self)
        dialog.setWindowTitle("🔍 찾은 파일 수")
        dialog.setGeometry(200, 200, 800, 600)

        layout = QVBoxLayout(dialog)

        # 통계 정보
        info_label = QLabel(f"검색어를 찾은 파일 {len(self.read_files)}개 (누적)")
        layout.addWidget(info_label)

        # 읽은 파일 목록
        text_edit = QTextEdit()
        text_edit.setReadOnly(True)
        text_edit.setFont(QFont("Consolas", 9))

        read_text = "\n".join(self.read_files)
        text_edit.setText(read_text)
        layout.addWidget(text_edit)

        # 닫기 버튼
        close_btn = QPushButton("닫기")
        close_btn.clicked.connect(dialog.close)
        layout.addWidget(close_btn)

        dialog.exec_()

    def show_no_match_files(self):
        """검색어 미포함 파일 목록 보기"""
        if not self.no_match_files:
            QMessageBox.information(self, "검색어 미포함 파일 수", "검색어 미포함 파일이 없습니다.\n\n먼저 검색을 실행해주세요.")
            return

        # 검색어 미포함 파일 목록을 보여주는 윈도우
        from PyQt5.QtWidgets import QDialog
        dialog = QDialog(self)
        dialog.setWindowTitle("❌ 검색어 미포함 파일 수")
        dialog.setGeometry(200, 200, 800, 600)

        layout = QVBoxLayout(dialog)

        # 통계 정보
        info_label = QLabel(f"검색어를 찾지 못한 파일 {len(self.no_match_files)}개 (누적)")
        layout.addWidget(info_label)

        # 검색어 미포함 파일 목록
        text_edit = QTextEdit()
        text_edit.setReadOnly(True)
        text_edit.setFont(QFont("Consolas", 9))

        no_match_text = "\n".join(self.no_match_files)
        text_edit.setText(no_match_text)
        layout.addWidget(text_edit)

        # 닫기 버튼
        close_btn = QPushButton("닫기")
        close_btn.clicked.connect(dialog.close)
        layout.addWidget(close_btn)

        dialog.exec_()

    def show_skipped_files(self):
        """스킵된 대용량파일 목록 보기"""
        if not self.skipped_files:
            QMessageBox.information(self, "스킵된 대용량파일", "스킵된 파일이 없습니다.\n\n먼저 검색을 실행해주세요.")
            return

        # 스킵된 대용량파일 목록을 보여주는 윈도우
        from PyQt5.QtWidgets import QDialog
        dialog = QDialog(self)
        dialog.setWindowTitle("⏭️ 스킵된 대용량파일")
        dialog.setGeometry(200, 200, 800, 600)

        layout = QVBoxLayout(dialog)

        # 통계 정보
        info_label = QLabel(f"10MB 이상 대용량파일 {len(self.skipped_files)}개 (누적)")
        layout.addWidget(info_label)

        # 스킵된 파일 목록
        text_edit = QTextEdit()
        text_edit.setReadOnly(True)
        text_edit.setFont(QFont("Consolas", 9))

        skipped_text = "\n".join(self.skipped_files)
        text_edit.setText(skipped_text)
        layout.addWidget(text_edit)

        # 닫기 버튼
        close_btn = QPushButton("닫기")
        close_btn.clicked.connect(dialog.close)
        layout.addWidget(close_btn)

        dialog.exec_()

    def sort_table_by_match_count(self):
        """테이블을 검색 단어수로 내림차순 정렬"""
        # results를 match_count로 정렬
        sorted_results = sorted(self.results, key=lambda x: x.get('match_count', 0), reverse=True)

        # 테이블 초기화
        self.table.setRowCount(0)

        # 정렬된 결과를 테이블에 다시 추가 (강조 표기 포함)
        shown_files = set()
        for result in sorted_results:
            file_path = result['full_path']
            if file_path not in shown_files:
                shown_files.add(file_path)

                # add_result_row와 동일한 로직으로 추가
                current_row_count = self.table.rowCount()
                self.table.insertRow(current_row_count)

                # 파일 크기 포맷
                size = result['size']
                if size < 1024:
                    size_str = f"{size} B"
                elif size < 1024 * 1024:
                    size_str = f"{size / 1024:.1f} KB"
                else:
                    size_str = f"{size / (1024 * 1024):.1f} MB"

                # 각 컬럼에 데이터 입력
                # 0: 파일명 (윈도우 탐색기 아이콘 + 검색어 강조 표시)
                keyword = getattr(self, 'current_keyword', '')
                use_regex = getattr(self, 'current_regex', False)
                filename_html = self._build_filename_display(result, keyword, use_regex)
                filename_label = QLabel()
                filename_label.setText(filename_html)
                filename_label.setStyleSheet("padding: 3px;")
                filename_label.setToolTip(result['full_path'])
                filename_label.setAttribute(Qt.WA_TransparentForMouseEvents)  # 마우스 이벤트를 테이블로 통과시켜 호버/더블클릭이 정상 동작하도록 함
                self.table.setCellWidget(current_row_count, 0, filename_label)
                # 위젯이 마우스 이벤트를 가로채지 않으므로, 툴팁/데이터는 동일 셀의 아이템에도 보관(다른 컬럼과 동일한 방식)
                filename_item = QTableWidgetItem("")
                filename_item.setToolTip(result['full_path'])
                filename_item.setData(Qt.UserRole, result['full_path'])
                filename_item.setData(Qt.UserRole + 1, result.get('matched_lines', []))
                self.table.setItem(current_row_count, 0, filename_item)

                # 1: 경로 (검색어 강조 표시)
                path_text = result['folder_path']
                if hasattr(self, 'current_keyword') and self.current_keyword:
                    highlighted_path = self._highlight_keyword(path_text, self.current_keyword, self.current_regex)
                    if '<span' in highlighted_path:
                        path_label = QLabel()
                        path_label.setText(highlighted_path)
                        path_label.setStyleSheet("padding: 3px;")
                        path_label.setToolTip(path_text)
                        self.table.setCellWidget(current_row_count, 1, path_label)
                    else:
                        path_item = QTableWidgetItem(path_text)
                        path_item.setToolTip(path_text)
                        path_item.setTextAlignment(Qt.AlignLeft | Qt.AlignVCenter)
                        self.table.setItem(current_row_count, 1, path_item)
                else:
                    path_item = QTableWidgetItem(path_text)
                    path_item.setToolTip(path_text)
                    path_item.setTextAlignment(Qt.AlignLeft | Qt.AlignVCenter)
                    self.table.setItem(current_row_count, 1, path_item)

                # 2: 크기
                size_item = QTableWidgetItem(size_str)
                size_item.setToolTip(str(result['size']) + " bytes")
                size_item.setTextAlignment(Qt.AlignRight | Qt.AlignVCenter)
                self.table.setItem(current_row_count, 2, size_item)

                # 3: 수정 날짜
                modified_item = QTableWidgetItem(result['modified'])
                modified_item.setTextAlignment(Qt.AlignLeft | Qt.AlignVCenter)
                self.table.setItem(current_row_count, 3, modified_item)

                # 4: 검색 단어수
                match_count = result.get('match_count', 0)
                match_count_item = QTableWidgetItem(str(match_count))
                match_count_item.setTextAlignment(Qt.AlignRight | Qt.AlignVCenter)
                self.table.setItem(current_row_count, 4, match_count_item)


def main():
    app = QApplication(sys.argv)
    window = FSearchGUI()
    window.show()
    sys.exit(app.exec_())


if __name__ == '__main__':
    multiprocessing.freeze_support()
    main()
